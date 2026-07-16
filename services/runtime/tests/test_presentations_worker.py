from __future__ import annotations

import asyncio
import hashlib
import json
import os
import subprocess
import sys
import uuid
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from local_host.plugins.executor import ManagedWorkerActionExecutor
from local_host.plugins.macos_vm import load_macos_vm_resources
from local_host.plugins.platforms import (
    current_managed_worker_execution_platform,
    current_managed_worker_platform,
)
from local_host.plugins.runtime_assets import RuntimeAssetHandle

REPO_ROOT = Path(__file__).resolve().parents[3]
WORKER = REPO_ROOT / "plugins" / "office" / "presentations" / "worker" / "presentations_worker.py"
GOLDEN_ROOT = REPO_ROOT / "plugins" / "office" / "presentations" / "golden"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _worker_command() -> tuple[str, ...]:
    frozen = os.environ.get("SHEJANE_TEST_PRESENTATIONS_WORKER")
    return (frozen,) if frozen else (sys.executable, str(WORKER))


def _invocation(
    action_id: str,
    arguments: dict[str, object],
    sources: list[tuple[str, Path, str]],
) -> dict[str, object]:
    return {
        "schema_version": 1,
        "invocation_id": str(uuid.uuid4()),
        "operation_id": f"run_01:{action_id}:001",
        "action": {
            "plugin_id": "org.shejane.presentations",
            "plugin_digest": "sha256:" + "f" * 64,
            "action_id": action_id,
        },
        "arguments": arguments,
        "inputs": [
            {
                "id": input_id,
                "path": f"/input/{path.name}",
                "media_type": media_type,
                "size_bytes": path.stat().st_size,
                "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            }
            for input_id, path, media_type in sources
        ],
        "grants": {"capabilities": []},
        "limits": {"timeout_ms": 15_000, "memory_mb": 512, "output_mb": 64},
    }


async def _invoke(
    tmp_path: Path,
    action_id: str,
    arguments: dict[str, object],
    *,
    sources: list[tuple[str, Path, str]] | None = None,
    runtime_assets: tuple[RuntimeAssetHandle, ...] = (),
) -> tuple[dict[str, object], Path]:
    roots = tmp_path / str(uuid.uuid4())
    input_root = roots / "input"
    output_root = roots / "output"
    input_root.mkdir(parents=True)
    output_root.mkdir()
    materialized: list[tuple[str, Path, str]] = []
    for input_id, source, media_type in sources or []:
        target = input_root / source.name
        target.write_bytes(source.read_bytes())
        materialized.append((input_id, target, media_type))
    vm_manifest = os.environ.get("SHEJANE_TEST_MACOS_VM_ASSETS")
    command = _worker_command()
    result = await ManagedWorkerActionExecutor(
        command,
        runtime_assets=runtime_assets,
        vm_resources=(
            load_macos_vm_resources(Path(vm_manifest).resolve(strict=True)) if vm_manifest else None
        ),
        package_root=Path(command[0]).resolve(strict=True).parent if vm_manifest else None,
    ).invoke(
        _invocation(action_id, arguments, materialized),
        input_root=input_root,
        output_root=output_root,
    )
    return result, output_root


def _tiny_png(path: Path) -> None:
    Image.new("RGB", (16, 16), (180, 30, 35)).save(path)


@pytest.mark.asyncio
async def test_presentations_worker_create_read_edit_and_preserve_input(tmp_path: Path) -> None:
    image = tmp_path / "pixel.png"
    _tiny_png(image)
    create_arguments = {
        "filename": "briefing.pptx",
        "slides": [
            {
                "layout": "Title Slide",
                "title": "Quarterly briefing",
                "subtitle": "SheJane",
                "notes": "Open with the outcome.",
            },
            {
                "layout": "Title and Content",
                "title": "Evidence",
                "bullets": [
                    {"text": "Revenue stable", "level": 0},
                    {"text": "Risk bounded", "level": 1},
                ],
                "notes": "Pause after the chart.",
                "text_boxes": [
                    {
                        "text": "CONFIDENTIAL",
                        "left": 9.8,
                        "top": 0.2,
                        "width": 3.0,
                        "height": 0.4,
                        "font_size": 10,
                    }
                ],
                "tables": [
                    {
                        "rows": [["Region", "Value"], ["North", "42"]],
                        "left": 0.8,
                        "top": 4.6,
                        "width": 5.0,
                        "height": 1.2,
                    }
                ],
                "images": [{"input_id": "logo", "left": 10.8, "top": 5.6, "width": 1.0}],
            },
        ],
    }
    created, create_output = await _invoke(
        tmp_path,
        "presentation.create",
        create_arguments,
        sources=[("logo", image, "image/png")],
    )
    replay, replay_output = await _invoke(
        tmp_path,
        "presentation.create",
        create_arguments,
        sources=[("logo", image, "image/png")],
    )
    created_path = create_output / "briefing.pptx"
    assert created["status"] == "succeeded", created
    assert replay["status"] == "succeeded", replay
    assert created_path.read_bytes() == (replay_output / "briefing.pptx").read_bytes()
    source_digest = hashlib.sha256(created_path.read_bytes()).hexdigest()

    read, _ = await _invoke(
        tmp_path,
        "presentation.read",
        {"input_id": "deck", "start_slide": 0, "max_slides": 10},
        sources=[("deck", created_path, PPTX_MIME)],
    )
    assert read["status"] == "succeeded"
    output = read["output"]
    assert output["slide_count"] == 2
    assert output["slides"][0]["title"] == "Quarterly briefing"
    assert output["slides"][0]["notes"] == "Open with the outcome."
    assert output["slides"][1]["bullets"] == [
        {"text": "Revenue stable", "level": 0},
        {"text": "Risk bounded", "level": 1},
    ]
    assert output["slides"][1]["image_count"] == 1
    assert output["slides"][1]["tables"][0]["rows"] == [
        ["Region", "Value"],
        ["North", "42"],
    ]
    assert output["unsupported_features"] == []

    edited, edit_output = await _invoke(
        tmp_path,
        "presentation.edit",
        {
            "input_id": "deck",
            "output_filename": "briefing-edited.pptx",
            "operations": [
                {
                    "type": "update_slide",
                    "index": 1,
                    "content": {
                        "title": "Verified evidence",
                        "bullets": [{"text": "Revenue +4%", "level": 0}],
                        "notes": "State the source.",
                    },
                },
                {
                    "type": "add_slide",
                    "slide": {"layout": "Title and Content", "title": "Next step", "bullets": []},
                },
                {"type": "move_slide", "from_index": 2, "to_index": 1},
            ],
        },
        sources=[("deck", created_path, PPTX_MIME)],
    )
    assert edited["status"] == "succeeded"
    assert hashlib.sha256(created_path.read_bytes()).hexdigest() == source_digest
    deck = Presentation(edit_output / "briefing-edited.pptx")
    assert [slide.shapes.title.text for slide in deck.slides] == [
        "Quarterly briefing",
        "Next step",
        "Verified evidence",
    ]
    assert "State the source." in deck.slides[2].notes_slide.notes_text_frame.text


@pytest.mark.asyncio
async def test_presentations_worker_returns_stable_input_errors(tmp_path: Path) -> None:
    for name, data, code in (
        ("broken.pptx", b"not a zip", "presentation_corrupt"),
        ("legacy.pptx", bytes.fromhex("d0cf11e0a1b11ae1"), "presentation_encrypted_or_legacy"),
    ):
        source = tmp_path / name
        source.write_bytes(data)
        result, _ = await _invoke(
            tmp_path,
            "presentation.read",
            {"input_id": "deck"},
            sources=[("deck", source, PPTX_MIME)],
        )
        assert result["status"] == "failed"
        assert result["error"]["code"] == code  # type: ignore[index]


@pytest.mark.asyncio
async def test_presentations_worker_renders_only_with_exact_runtime_asset(tmp_path: Path) -> None:
    created, created_output = await _invoke(
        tmp_path,
        "presentation.create",
        {"filename": "render.pptx", "slides": [{"layout": "Title Slide", "title": "Rendered"}]},
    )
    assert created["status"] == "succeeded"
    source = created_output / "render.pptx"

    asset_root = tmp_path / "office-runtime"
    asset_root.mkdir()
    soffice = asset_root / "soffice"
    mutool = asset_root / "mutool"
    soffice.write_text(
        f"#!{sys.executable}\n"
        "import pathlib, sys\n"
        "args = sys.argv[1:]\n"
        "out = pathlib.Path(args[args.index('--outdir') + 1])\n"
        "source = pathlib.Path(args[-1])\n"
        "(out / (source.stem + '.pdf')).write_bytes(b'%PDF-1.7\\n%%EOF\\n')\n",
        encoding="utf-8",
    )
    mutool.write_text(
        f"#!{sys.executable}\n"
        "import pathlib, sys\n"
        "args = sys.argv[1:]\n"
        "if args[0] == 'info':\n"
        "    print('Pages: 1')\n"
        "else:\n"
        "    pathlib.Path(args[args.index('-o') + 1] % 1).write_bytes(b'PNG')\n",
        encoding="utf-8",
    )
    soffice.chmod(0o500)
    mutool.chmod(0o500)
    (asset_root / "office-runtime.json").write_text(
        json.dumps({"soffice": "soffice", "mutool": "mutool"}), encoding="utf-8"
    )
    sbom = asset_root / "sbom.spdx.json"
    sbom.write_text("{}", encoding="utf-8")
    asset = RuntimeAssetHandle(
        asset_id="org.libreoffice.runtime",
        version="25.8.7",
        platform="darwin/arm64",
        digest="sha256:" + "a" * 64,
        root=asset_root,
        payload=asset_root,
        license="MPL-2.0",
        source_url="https://www.libreoffice.org/",
        sbom=sbom,
    )

    rendered, output = await _invoke(
        tmp_path,
        "presentation.render",
        {"input_id": "deck", "include_png": True, "max_slides": 1, "dpi": 144},
        sources=[("deck", source, PPTX_MIME)],
        runtime_assets=(asset,),
    )
    assert rendered["status"] == "succeeded"
    assert rendered["output"] == {
        "slide_count": 1,
        "rendered_slides": 1,
        "pdf_name": "render.pdf",
        "png_names": ["render.slide-0001.png"],
        "unsupported_features": [],
        "warnings": [],
    }
    assert {path.name for path in output.iterdir()} == {"render.pdf", "render.slide-0001.png"}


@pytest.mark.asyncio
async def test_presentations_worker_real_runtime_render_golden(tmp_path: Path) -> None:
    configured = os.environ.get("SHEJANE_TEST_OFFICE_RUNTIME")
    if not configured:
        pytest.skip("real Office runtime asset not configured")
    asset_root = Path(configured).resolve(strict=True)
    target = (
        current_managed_worker_execution_platform()
        if os.environ.get("SHEJANE_TEST_MACOS_VM_ASSETS")
        else current_managed_worker_platform()
    )
    assert target is not None
    asset = RuntimeAssetHandle(
        asset_id="org.libreoffice.runtime",
        version="25.8.7",
        platform=target,
        digest="sha256:" + "a" * 64,
        root=asset_root,
        payload=asset_root / "payload",
        license="MPL-2.0 AND AGPL-3.0-only AND OFL-1.1",
        source_url="https://www.libreoffice.org/",
        sbom=asset_root / "payload" / "office-runtime.json",
    )
    image = tmp_path / "mark.png"
    Image.new("RGB", (160, 160), (180, 30, 35)).save(image)
    created, created_output = await _invoke(
        tmp_path,
        "presentation.create",
        {
            "filename": "golden.pptx",
            "slides": [
                {
                    "layout": "Title Slide",
                    "title": "石间演示文稿",
                    "subtitle": "日本語 · 한국어",
                    "notes": "Opening note",
                    "font_name": "Noto Sans CJK SC",
                },
                {
                    "layout": "Title and Content",
                    "title": "Evidence",
                    "font_name": "Noto Sans CJK SC",
                    "bullets": [
                        {"text": "可验证的结果", "level": 0, "font_name": "Noto Sans CJK SC"},
                        {"text": "Preserved layout", "level": 1, "font_name": "Noto Sans CJK SC"},
                    ],
                    "notes": "Second slide note",
                    "text_boxes": [
                        {
                            "text": "Noto Sans CJK",
                            "left": 8.8,
                            "top": 0.3,
                            "width": 3.8,
                            "height": 0.5,
                            "font_name": "Noto Sans CJK SC",
                            "font_size": 14,
                        }
                    ],
                    "tables": [
                        {
                            "rows": [["Metric", "Value"], ["Quality", "Stable"]],
                            "left": 0.8,
                            "top": 4.5,
                            "width": 5.4,
                            "height": 1.3,
                        }
                    ],
                    "images": [{"input_id": "mark", "left": 10.8, "top": 5.2, "width": 1.2}],
                },
            ],
        },
        sources=[("mark", image, "image/png")],
    )
    assert created["status"] == "succeeded", created
    source = created_output / "golden.pptx"
    source_digest = hashlib.sha256(source.read_bytes()).hexdigest()
    rendered, output = await _invoke(
        tmp_path,
        "presentation.render",
        {"input_id": "deck", "include_png": True, "max_slides": 4, "dpi": 144},
        sources=[("deck", source, PPTX_MIME)],
        runtime_assets=(asset,),
    )
    assert rendered["status"] == "succeeded", rendered
    assert rendered["output"]["slide_count"] == 2  # type: ignore[index]
    assert rendered["output"]["unsupported_features"] == []  # type: ignore[index]
    assert hashlib.sha256(source.read_bytes()).hexdigest() == source_digest

    if not os.environ.get("SHEJANE_TEST_MACOS_VM_ASSETS"):
        config = json.loads((asset.payload / "office-runtime.json").read_text(encoding="utf-8"))
        mutool = asset.payload / config["mutool"]
        extracted = (
            await asyncio.to_thread(
                subprocess.run,
                [str(mutool), "draw", "-F", "txt", "-o", "-", str(output / "golden.pdf")],
                check=True,
                capture_output=True,
                text=True,
            )
        ).stdout
        assert "石间演示文稿" in extracted
        assert "Preserved layout" in extracted

    expected = json.loads((GOLDEN_ROOT / f"{target.replace('/', '-')}.json").read_text())
    actual = []
    for page in range(1, rendered["output"]["rendered_slides"] + 1):  # type: ignore[index,operator]
        with Image.open(output / f"golden.slide-{page:04d}.png") as rendered_image:
            gray = rendered_image.convert("L")
            actual.append(
                {
                    "size": list(gray.size),
                    "ink_bbox": list(_ink_bbox(gray)),
                    "dhash": _dhash(gray),
                }
            )
    assert len(actual) == expected["slide_count"]
    assert actual == expected["slides"]


def _ink_bbox(image: Image.Image) -> tuple[int, int, int, int]:
    return image.point(lambda value: 255 - value).getbbox() or (0, 0, 0, 0)


def _dhash(image: Image.Image) -> str:
    sampled = image.resize((9, 8), Image.Resampling.LANCZOS)
    pixels = list(sampled.get_flattened_data())
    value = 0
    for row in range(8):
        for column in range(8):
            value = (value << 1) | (pixels[row * 9 + column] > pixels[row * 9 + column + 1])
    return f"{value:016x}"
