"""office.* — read, outline, and edit tools for Word (.docx) and Excel (.xlsx).

Phase 1 (already shipped) — read-only:
  * office.read       — full content as LLM-ready markdown
  * office.outline    — cheap structural summary

Phase 2 (this module) — edits via COPY-ON-FIRST-WRITE:
  Docx: find_replace, insert_paragraph, update_paragraph,
        delete_paragraph, apply_style
  Xlsx: set_cells, set_formula, set_cell_format, merge_cells,
        add_row, read_range

The user's hard constraint for Phase 2: the original file is NEVER
modified. Every write tool copies the original to a sibling
`<basename>.edited.<ext>` on first write and targets the copy
thereafter. Repeated edits land in the same copy. The user can reset
edits by deleting `xxx.edited.docx` in Finder.

Because the original is untouched, write tools do NOT require HITL
approval (they're not registered in `agent/builder.py:DESTRUCTIVE_TOOLS`).
This trades a permission prompt for the predictable "copy is the
safety net" contract.

All write tools also use an atomic write pattern: write to
`<target>.tmp`, re-open with the appropriate library to verify it's a
valid OOXML file, then `os.replace(tmp, target)`. A mid-write failure
leaves `target` exactly as it was (the last known-good edit).

Path safety: the daemon trusts the agent to pass paths that came out
of `workspace.open`. We do extension + existence checks here; the
fs middleware enforces the workspace-membership rule for fs.* tools.
"""

from __future__ import annotations

import os
import re
import shutil
import tempfile
from collections.abc import Callable
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from docx.document import Document as DocxDocumentType
from docx.text.paragraph import Paragraph as DocxParagraph
from langchain_core.tools import tool
from markitdown import MarkItDown
from openpyxl import load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import column_index_from_string, get_column_letter
from openpyxl.utils.cell import range_boundaries
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

# Module-level converter — MarkItDown is cheap to construct but creating it
# once is even cheaper, and it's thread-safe for the read API we use.
_md = MarkItDown()

# Hard cap on returned markdown size. A 50-page docx is typically ~30 KB
# of text; a 100-sheet xlsx ingested as markdown tables can blow past
# 200 KB. We truncate so the LLM context doesn't get nuked by one
# `office.read` call. The agent gets told upfront in the return shape.
_MARKDOWN_CHAR_CAP = 60_000

# Mapping from extension to the human-readable kind we surface to the
# client. Used by both office.read and office.outline so the frontend
# can pick the right preview component.
_EXTENSION_KIND = {
    ".docx": "word",
    ".xlsx": "excel",
    ".pptx": "powerpoint",
}

# Naming for editable copies. The infix sits between basename and
# extension: `report.docx` → `report.edited.docx`. Hard-coded because
# the frontend + developer.md both reference the convention.
_EDITED_INFIX = "edited"


# ═══════════════════════════════════════════════════════════════════════
# Path validation + write-target resolution
# ═══════════════════════════════════════════════════════════════════════


def _validate_path(path: str) -> tuple[str | None, str | None, str | None]:
    """Resolve and validate the file path.

    Returns (resolved_path, kind, error). On success error is None; on
    failure resolved_path/kind are None and error describes what went wrong.
    """
    if not path:
        return None, None, "path required"
    resolved = os.path.abspath(os.path.expanduser(path))
    if not os.path.isfile(resolved):
        return None, None, f"file not found: {resolved}"
    ext = Path(resolved).suffix.lower()
    kind = _EXTENSION_KIND.get(ext)
    if not kind:
        return (
            None,
            None,
            (
                f"unsupported extension {ext!r}; office.* only handles .docx and .xlsx "
                "(use read_file for plain text, image.* for images)"
            ),
        )
    return resolved, kind, None


def _resolve_write_target(original_path: str) -> str:
    """Return the path of the editable copy for `original_path`.

    Convention: `<basename>.edited.<ext>` in the same directory.
    Calling on an already-edited file returns it unchanged (idempotent),
    so chained edits keep targeting the same copy.
    """
    p = Path(original_path)
    if p.stem.endswith(f".{_EDITED_INFIX}"):
        return str(p)
    return str(p.with_name(f"{p.stem}.{_EDITED_INFIX}{p.suffix}"))


def _ensure_copy_for_write(original_path: str) -> str:
    """Materialize the editable copy if it doesn't exist; return its path.

    Uses `shutil.copy2` so the copy starts with the original's mtime
    and permission bits — the user can sort by date in Finder and see
    the edit relative to the original.
    """
    target = _resolve_write_target(original_path)
    if target == original_path:
        return target
    if not os.path.exists(target):
        shutil.copy2(original_path, target)
    return target


def _verify_file(path: str, kind: str) -> None:
    """Open a freshly written file with the relevant library; any
    exception means the write produced an invalid OOXML file and the
    caller must NOT promote it to the target."""
    if kind == "word":
        DocxDocument(path)
    elif kind == "excel":
        wb = load_workbook(path, read_only=True)
        wb.close()
    elif kind == "powerpoint":
        # python-pptx raises on bad XML / bad ZIP just like python-docx
        # does for .docx files. No explicit close needed.
        PptxPresentation(path)


def _atomic_write(
    target: str,
    kind: str,
    write_fn: Callable[[str], dict[str, Any]],
) -> dict[str, Any]:
    """Write to a sibling tmp file → verify → atomically replace `target`.

    `write_fn(tmp_path)` performs the actual save and returns a dict
    used as the `summary` field of the tool result. If write_fn or the
    verification raises, the tmp file is cleaned up and `target` keeps
    its previous content.

    Tmp naming: hidden file in the same directory as target, with the
    SAME extension as target (`.docx` / `.xlsx`). openpyxl dispatches
    its save format from the filename extension and rejects unknown
    suffixes (including `.tmp`), so we can't just append `.tmp`.
    """
    target_path = Path(target)
    fd, tmp = tempfile.mkstemp(
        suffix=target_path.suffix,
        prefix=f".{target_path.stem}.tmp.",
        dir=str(target_path.parent),
    )
    # mkstemp opens the file for us; we only want the unique path.
    os.close(fd)
    try:
        summary = write_fn(tmp)
        _verify_file(tmp, kind)
        os.replace(tmp, target)
        return summary
    except Exception:
        if os.path.exists(tmp):
            try:
                os.remove(tmp)
            except OSError:
                pass
        raise


def _write_result(
    original: str,
    target: str,
    kind: str,
    summary: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Build the standard success-response dict every write tool returns.

    Keys (LLM-visible contract):
      ok            "true"
      original_path the path the user passed
      target_path   alias retained for forwards-compat (= edited_path)
      edited_path   absolute path of the copy that now holds the edits
      kind          "word" or "excel"
      summary       per-tool details (counts, ranges affected, etc.)
    """
    return {
        "ok": "true",
        "original_path": original,
        "edited_path": target,
        "kind": kind,
        "summary": summary or {},
    }


def _write_error(
    kind: str | None,
    message: str,
    original: str | None = None,
) -> dict[str, Any]:
    """Build the standard failure-response dict every write tool returns."""
    out: dict[str, Any] = {"ok": "false", "error": message}
    if kind:
        out["kind"] = kind
    if original:
        out["original_path"] = original
    return out


# ═══════════════════════════════════════════════════════════════════════
# READ tools (Phase 1)
# ═══════════════════════════════════════════════════════════════════════


@tool("office.read")
def office_read(path: str) -> dict[str, Any]:
    """Read a Word (.docx) or Excel (.xlsx) file as LLM-ready markdown.

    Prefer this over `read_file` for these extensions — `read_file` would
    return the raw ZIP/XML which is useless for analysis. This tool runs
    `markitdown` which converts headings, paragraphs, tables, and cells
    into clean markdown the LLM can reason about directly.

    Does NOT open the right-side document preview panel. If you want the
    user to see the file rendered, mention the filename in your reply
    so the renderer makes it clickable. The preview opens when the user
    clicks, or after a successful office.* WRITE tool completes.

    Args:
        path: Absolute filesystem path to a .docx or .xlsx file. Must
              already exist. Use `workspace.open` first if the file
              lives in a directory the agent hasn't been authorized for.

    Returns:
        dict with keys:
          ok ("true" / "false")
          path (echoed back, absolute)
          kind ("word" or "excel")
          markdown (the converted markdown content, possibly truncated)
          truncated ("true" / "false") — set when content exceeded the cap
          error (only present when ok="false")
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return {"ok": "false", "error": err}
    assert resolved is not None and kind is not None  # for type checker
    try:
        result = _md.convert(resolved)
    except Exception as exc:
        return {
            "ok": "false",
            "path": resolved,
            "kind": kind,
            "error": f"failed to convert {kind}: {exc.__class__.__name__}: {exc}",
        }
    text = result.text_content or ""
    truncated = len(text) > _MARKDOWN_CHAR_CAP
    if truncated:
        text = (
            text[:_MARKDOWN_CHAR_CAP]
            + "\n\n…(truncated, full size = "
            + str(len(result.text_content))
            + " chars)"
        )
    return {
        "ok": "true",
        "path": resolved,
        "kind": kind,
        "markdown": text,
        "truncated": "true" if truncated else "false",
    }


def _outline_docx(path: str) -> dict[str, Any]:
    """Cheap structural summary of a .docx — heading text and paragraph count."""
    doc = DocxDocument(path)
    headings: list[dict[str, Any]] = []
    paragraph_count = 0
    table_count = len(doc.tables)
    for para in doc.paragraphs:
        paragraph_count += 1
        style_name = (para.style.name if para.style is not None else "") or ""
        if style_name.startswith("Heading"):
            level_str = style_name.split()[-1] if " " in style_name else "1"
            try:
                level = int(level_str)
            except ValueError:
                level = 1
            text = para.text.strip()
            if text:
                headings.append({"level": level, "text": text})
    return {
        "headings": headings,
        "paragraph_count": paragraph_count,
        "table_count": table_count,
    }


def _outline_xlsx(path: str) -> dict[str, Any]:
    """Cheap structural summary of a .xlsx — sheet names + dimensions."""
    wb = load_workbook(path, read_only=True, data_only=True)
    sheets: list[dict[str, Any]] = []
    try:
        for name in wb.sheetnames:
            ws = wb[name]
            sheets.append(
                {
                    "name": name,
                    "rows": ws.max_row or 0,
                    "columns": ws.max_column or 0,
                }
            )
    finally:
        wb.close()
    return {"sheets": sheets}


def _slide_title_text(slide) -> str:
    """Extract the title text from a slide, or "" if no title placeholder."""
    # python-pptx exposes a `.shapes.title` shortcut, but it's None
    # when the slide layout has no title placeholder. Fall back to the
    # first text-bearing shape so we still surface SOMETHING for
    # title-less slides (cover layouts use centered text frames).
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        text = slide.shapes.title.text_frame.text.strip()
        if text:
            return text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text.splitlines()[0]
    return ""


def _slide_bullets(slide) -> list[str]:
    """Extract paragraph-level text from the slide's content placeholder.

    A slide typically has one title placeholder + one body placeholder.
    We pull text from every non-title shape's text_frame and split by
    paragraph. Empty paragraphs are dropped.
    """
    bullets: list[str] = []
    title_shape = slide.shapes.title
    # python-pptx returns different wrapper instances for the same
    # underlying XML element across `shapes.title` vs iteration, so an
    # `is` check fails. Compare the underlying lxml elements instead.
    title_elem = title_shape._element if title_shape is not None else None
    for shape in slide.shapes:
        if title_elem is not None and shape._element is title_elem:
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                bullets.append(text)
    return bullets


def _slide_notes(slide) -> str:
    """Speaker notes for a slide; "" when no notes slide is attached."""
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def _outline_pptx(path: str) -> dict[str, Any]:
    """Structural summary of a .pptx — per-slide title, bullets, notes,
    shape counts. Both `office.outline` and `office.read_slides` consume
    this; the panel preview also fetches it through a small HTTP
    endpoint."""
    prs = PptxPresentation(path)
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        slides.append(
            {
                "index": idx,
                "layout": slide.slide_layout.name if slide.slide_layout else "",
                "title": _slide_title_text(slide),
                "bullets": _slide_bullets(slide),
                "notes": _slide_notes(slide),
                "shape_count": len(slide.shapes),
                "image_count": sum(1 for sh in slide.shapes if sh.shape_type == 13),  # 13 = PICTURE
            }
        )
    return {"slides": slides, "slide_count": len(slides)}


@tool("office.outline")
def office_outline(path: str) -> dict[str, Any]:
    """Return a cheap structural summary of a .docx, .xlsx, or .pptx file.

    Use this BEFORE `office.read` when the file is large and you only need
    to know what's in it — for example: "tell me what sheets are in
    Q4.xlsx", "does report.docx have a section about pricing?", or "how
    many slides does pitch.pptx have?". Reading the outline is
    O(metadata); reading the full markdown is O(file).

    Args:
        path: Absolute filesystem path to a .docx, .xlsx, or .pptx file.

    Returns:
        dict with keys:
          ok ("true" / "false")
          path (echoed back, absolute)
          kind ("word", "excel", or "powerpoint")
          For .docx: headings, paragraph_count, table_count.
          For .xlsx: sheets (list of {name, rows, columns}).
          For .pptx: slides (list of {index, layout, title, bullets,
                     notes, shape_count, image_count}), slide_count.
          error (only present when ok="false")
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return {"ok": "false", "error": err}
    assert resolved is not None and kind is not None
    try:
        if kind == "word":
            details = _outline_docx(resolved)
        elif kind == "excel":
            details = _outline_xlsx(resolved)
        else:  # powerpoint
            details = _outline_pptx(resolved)
    except Exception as exc:
        return {
            "ok": "false",
            "path": resolved,
            "kind": kind,
            "error": f"failed to read {kind} outline: {exc.__class__.__name__}: {exc}",
        }
    return {
        "ok": "true",
        "path": resolved,
        "kind": kind,
        **details,
    }


# ═══════════════════════════════════════════════════════════════════════
# DOCX write helpers
# ═══════════════════════════════════════════════════════════════════════


def _find_paragraph(doc: DocxDocumentType, target: str) -> DocxParagraph | None:
    """Return the first paragraph whose .text contains `target` (exact
    substring match), or None. Tables are walked too — paragraphs inside
    cells should be addressable just like top-level body paragraphs."""
    if not target:
        return None
    for para in doc.paragraphs:
        if target in para.text:
            return para
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if target in para.text:
                        return para
    return None


def _set_paragraph_text(para: DocxParagraph, new_text: str) -> None:
    """Replace a paragraph's text with `new_text`, collapsing all
    existing runs into one. This LOSES run-level formatting inside the
    paragraph (bold-in-middle, font color, etc.) but preserves the
    paragraph-level style (Heading 1, Quote, list bullet, etc.). For
    the find_replace / update_paragraph use cases this is the right
    trade — agents are editing content, not micro-formatting."""
    if para.runs:
        # Reuse the first run to keep its char formatting as the new
        # paragraph's font (best-effort), then clear the rest.
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.add_run(new_text)


def _insert_paragraph_relative(
    anchor: DocxParagraph, content: str, position: str, style: str | None
) -> DocxParagraph:
    """Insert a new paragraph immediately before or after `anchor`.

    python-docx has no built-in "insert before/after" — we manipulate
    the XML directly. This is the standard recipe from the docx-py
    cookbook.
    """
    from copy import deepcopy

    new_p = deepcopy(anchor._p)
    # Strip the existing runs from the cloned XML so we get a blank
    # paragraph element with the same pPr (paragraph properties).
    for child in list(new_p):
        if child.tag.endswith("}r"):
            new_p.remove(child)
    if position == "before":
        anchor._p.addprevious(new_p)
    else:  # default + 'after'
        anchor._p.addnext(new_p)
    new_para = DocxParagraph(new_p, anchor._parent)
    if style:
        new_para.style = style
    new_para.add_run(content)
    return new_para


# ═══════════════════════════════════════════════════════════════════════
# DOCX write tools
# ═══════════════════════════════════════════════════════════════════════


@tool("office.find_replace")
def office_find_replace(
    path: str,
    find: str,
    replace: str,
    count: int | None = None,
) -> dict[str, Any]:
    """Replace text across a .docx. Writes to a copy named `<basename>.edited.docx`.

    The original file is NEVER modified — you always operate on a copy.
    If the copy already exists from a previous edit, this writes to
    that same copy (chained edits land in one file).

    Args:
        path: .docx to edit. May be the original or its `.edited` copy
              (idempotent: if you pass the edited copy, we write back
              to it).
        find: text to search for (exact substring match).
        replace: text to substitute in.
        count: optional — stop after this many replacements. None =
               replace all.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={replaced}`. Use `edited_path` for subsequent reads
        or edits on this document.

    Note on formatting: this tool collapses run-level char formatting
    in each modified paragraph (paragraph-level style is preserved).
    If you need to keep "bold middle word" intact, use
    `office.update_paragraph` instead — it's whole-paragraph anyway.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "word":
        return _write_error(kind, "office.find_replace requires a .docx file", resolved)
    if not find:
        return _write_error(kind, "find text required", resolved)

    target = _ensure_copy_for_write(resolved)
    remaining = count if (count is None or count > 0) else 0

    def _do_write(tmp_path: str) -> dict[str, Any]:
        nonlocal remaining
        doc = DocxDocument(target)
        replaced = 0
        for para in _iter_all_paragraphs(doc):
            if remaining == 0:
                break
            text = para.text
            if find not in text:
                continue
            occurrences = text.count(find)
            if remaining is not None:
                occurrences = min(occurrences, remaining)
            new_text = text.replace(find, replace, occurrences)
            _set_paragraph_text(para, new_text)
            replaced += occurrences
            if remaining is not None:
                remaining -= occurrences
        doc.save(tmp_path)
        return {"replaced": replaced}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


def _iter_all_paragraphs(doc: DocxDocumentType):
    """Yield body paragraphs + paragraphs nested in table cells."""
    yield from doc.paragraphs
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs


@tool("office.insert_paragraph")
def office_insert_paragraph(
    path: str,
    anchor: str,
    content: str,
    position: str = "after",
    style: str | None = None,
) -> dict[str, Any]:
    """Insert a new paragraph relative to an anchor paragraph in a .docx.

    Writes to `<basename>.edited.docx`; original preserved.

    Args:
        path: .docx to edit.
        anchor: text that uniquely identifies the anchor paragraph
                (first paragraph whose text contains this substring).
                Empty string OR the literal "__END__" inserts at the
                very end of the document.
        content: text of the new paragraph.
        position: "before" or "after" the anchor (default "after").
                  Ignored when anchor is "__END__".
        style: optional python-docx style name (e.g. "Heading 1",
               "Normal", "Quote", "List Bullet"). Use `office.outline`
               to see which styles already appear in the document.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={inserted_at, style?}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "word":
        return _write_error(kind, "office.insert_paragraph requires a .docx file", resolved)
    if position not in {"before", "after"}:
        return _write_error(
            kind, f"position must be 'before' or 'after', got {position!r}", resolved
        )

    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        doc = DocxDocument(target)
        if anchor in ("", "__END__"):
            new_para = (
                doc.add_paragraph(content, style=style) if style else doc.add_paragraph(content)
            )
            doc.save(tmp_path)
            return {"inserted_at": "end", "style": new_para.style.name if new_para.style else None}
        match = _find_paragraph(doc, anchor)
        if match is None:
            raise ValueError(f"anchor not found: {anchor!r}")
        new_para = _insert_paragraph_relative(match, content, position, style)
        doc.save(tmp_path)
        return {
            "inserted_at": f"{position} anchor",
            "anchor_excerpt": match.text[:60],
            "style": new_para.style.name if new_para.style else None,
        }

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.update_paragraph")
def office_update_paragraph(
    path: str,
    target: str,
    content: str,
    style: str | None = None,
) -> dict[str, Any]:
    """Rewrite the first paragraph that contains `target` substring.

    Writes to `<basename>.edited.docx`; original preserved.

    Args:
        path: .docx to edit.
        target: substring that identifies the paragraph to update
                (first match wins; pick something unique).
        content: new full text of that paragraph (replaces ALL existing
                 text in the paragraph).
        style: optional new paragraph style name.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={updated_excerpt, style?}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "word":
        return _write_error(kind, "office.update_paragraph requires a .docx file", resolved)
    if not target:
        return _write_error(kind, "target text required", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        doc = DocxDocument(target_path)
        match = _find_paragraph(doc, target)
        if match is None:
            raise ValueError(f"target paragraph not found: {target!r}")
        old_excerpt = match.text[:60]
        _set_paragraph_text(match, content)
        if style:
            match.style = style
        doc.save(tmp_path)
        return {
            "updated_excerpt": old_excerpt,
            "style": match.style.name if match.style else None,
        }

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


@tool("office.delete_paragraph")
def office_delete_paragraph(path: str, target: str) -> dict[str, Any]:
    """Delete the first paragraph that contains `target` substring.

    Writes to `<basename>.edited.docx`; original preserved.

    Args:
        path: .docx to edit.
        target: substring that identifies the paragraph (first match wins).

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={deleted_excerpt}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "word":
        return _write_error(kind, "office.delete_paragraph requires a .docx file", resolved)
    if not target:
        return _write_error(kind, "target text required", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        doc = DocxDocument(target_path)
        match = _find_paragraph(doc, target)
        if match is None:
            raise ValueError(f"target paragraph not found: {target!r}")
        deleted_excerpt = match.text[:60]
        # Detach the paragraph element from its parent (python-docx has
        # no public `.delete()`; we manipulate XML directly).
        p_element = match._p
        p_element.getparent().remove(p_element)
        doc.save(tmp_path)
        return {"deleted_excerpt": deleted_excerpt}

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


@tool("office.apply_style")
def office_apply_style(path: str, target: str, style: str) -> dict[str, Any]:
    """Change the paragraph style of the first paragraph matching `target`.

    Writes to `<basename>.edited.docx`; original preserved.

    Args:
        path: .docx to edit.
        target: substring identifying the paragraph (first match wins).
        style: python-docx style name. Common values: "Heading 1",
               "Heading 2", "Heading 3", "Normal", "Quote", "Title",
               "Subtitle", "List Bullet", "List Number".

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={paragraph_excerpt, old_style, new_style}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "word":
        return _write_error(kind, "office.apply_style requires a .docx file", resolved)
    if not target or not style:
        return _write_error(kind, "target and style are both required", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        doc = DocxDocument(target_path)
        match = _find_paragraph(doc, target)
        if match is None:
            raise ValueError(f"target paragraph not found: {target!r}")
        old_style = match.style.name if match.style else None
        match.style = style
        doc.save(tmp_path)
        return {
            "paragraph_excerpt": match.text[:60],
            "old_style": old_style,
            "new_style": match.style.name if match.style else None,
        }

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


# ═══════════════════════════════════════════════════════════════════════
# XLSX write helpers
# ═══════════════════════════════════════════════════════════════════════


_CELL_REF_RE = re.compile(r"^([A-Za-z]+)(\d+)$")


def _parse_cell_ref(ref: str) -> tuple[int, int]:
    """Parse an A1-style cell ref into (row, col) 1-indexed tuple."""
    m = _CELL_REF_RE.match(ref.strip())
    if not m:
        raise ValueError(f"invalid cell reference: {ref!r} (expected like 'A1')")
    col_str, row_str = m.groups()
    return int(row_str), column_index_from_string(col_str.upper())


def _parse_range(range_ref: str) -> tuple[int, int, int, int]:
    """Parse a range like 'A1:C3' (or 'A1' for a single cell) into
    (min_row, min_col, max_row, max_col), all 1-indexed."""
    ref = range_ref.strip()
    if ":" not in ref:
        # Single cell — treat as a 1x1 range.
        row, col = _parse_cell_ref(ref)
        return row, col, row, col
    min_col, min_row, max_col, max_row = range_boundaries(ref.upper())
    return min_row, min_col, max_row, max_col


def _resolve_sheet(wb, sheet: str | None):
    """Return the worksheet by name, or the active sheet when sheet is None."""
    if not sheet:
        return wb.active
    if sheet not in wb.sheetnames:
        raise ValueError(f"sheet not found: {sheet!r}; available: {wb.sheetnames}")
    return wb[sheet]


def _color_to_hex(color: str) -> str:
    """Normalize a CSS-ish color string into openpyxl's 8-char AARRGGBB
    or 6-char RRGGBB hex (without leading '#'). Accepts '#RRGGBB',
    'RRGGBB', '#AARRGGBB', 'AARRGGBB'. Bare names not supported."""
    c = color.strip().lstrip("#").upper()
    if len(c) not in (6, 8) or not all(ch in "0123456789ABCDEF" for ch in c):
        raise ValueError(f"color must be hex like '#FF5722' or '#80FF5722', got {color!r}")
    return c


# ═══════════════════════════════════════════════════════════════════════
# XLSX write tools
# ═══════════════════════════════════════════════════════════════════════


@tool("office.set_cells")
def office_set_cells(
    path: str,
    sheet: str | None,
    range: str,
    values: list[list[Any]],
) -> dict[str, Any]:
    """Write a 2D values block into a rectangular cell range.

    Writes to `<basename>.edited.xlsx`; original preserved.

    Args:
        path: .xlsx to edit.
        sheet: sheet name (or None for the active sheet).
        range: A1-style range like "A1:C3" or single cell "B5".
        values: 2D list shaped to match the range. Top-left is range
                top-left. Each inner list is one row. Extra cells in
                the range that lack values stay unchanged; extra
                values beyond the range are ignored.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={range, rows_written, cols_written, cells_written}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "excel":
        return _write_error(kind, "office.set_cells requires a .xlsx file", resolved)
    if not values or not isinstance(values, list):
        return _write_error(kind, "values must be a non-empty 2D list", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        wb = load_workbook(target_path)
        try:
            ws = _resolve_sheet(wb, sheet)
            min_row, min_col, max_row, max_col = _parse_range(range)
            rng_rows = max_row - min_row + 1
            rng_cols = max_col - min_col + 1
            written = 0
            for r_offset, row_values in enumerate(values):
                if r_offset >= rng_rows:
                    break
                if not isinstance(row_values, list):
                    raise ValueError(f"values[{r_offset}] must be a list")
                for c_offset, cell_value in enumerate(row_values):
                    if c_offset >= rng_cols:
                        break
                    ws.cell(row=min_row + r_offset, column=min_col + c_offset, value=cell_value)
                    written += 1
            wb.save(tmp_path)
        finally:
            wb.close()
        return {
            "range": range,
            "rows_written": min(len(values), rng_rows),
            "cols_written": rng_cols,
            "cells_written": written,
        }

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


@tool("office.set_formula")
def office_set_formula(
    path: str,
    sheet: str | None,
    cell: str,
    formula: str,
) -> dict[str, Any]:
    """Write a formula into one cell. Writes to `<basename>.edited.xlsx`.

    NOTE: openpyxl writes the formula text — it does NOT evaluate it.
    The cell's displayed value updates only when Microsoft Excel /
    LibreOffice / Numbers opens the file. The right-side preview shows
    the literal formula text until then. This is an openpyxl limitation,
    not a bug in this tool.

    Args:
        path: .xlsx to edit.
        sheet: sheet name (None = active).
        cell: target cell like "D2".
        formula: must start with "=", e.g. "=SUM(A2:A10)",
                 "=IF(B2>0, \\"Pos\\", \\"Neg\\")", "=VLOOKUP(...)".

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={cell, formula}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "excel":
        return _write_error(kind, "office.set_formula requires a .xlsx file", resolved)
    formula = (formula or "").strip()
    if not formula.startswith("="):
        return _write_error(kind, "formula must start with '='", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        wb = load_workbook(target_path)
        try:
            ws = _resolve_sheet(wb, sheet)
            row, col = _parse_cell_ref(cell)
            ws.cell(row=row, column=col, value=formula)
            wb.save(tmp_path)
        finally:
            wb.close()
        return {"cell": cell, "formula": formula}

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


@tool("office.set_cell_format")
def office_set_cell_format(
    path: str,
    sheet: str | None,
    range: str,
    bold: bool | None = None,
    italic: bool | None = None,
    font_size: float | None = None,
    font_color: str | None = None,
    bg_color: str | None = None,
    align: str | None = None,
    border: bool | None = None,
) -> dict[str, Any]:
    """Apply font / fill / alignment / border formatting to a cell range.

    Writes to `<basename>.edited.xlsx`; original preserved. Only the
    arguments you pass are applied — None means "leave unchanged".

    Args:
        path: .xlsx to edit.
        sheet: sheet name (None = active).
        range: A1-style range.
        bold, italic: True / False (None to leave alone).
        font_size: e.g. 14.
        font_color, bg_color: hex like "#FF5722" or "FF5722". Alpha
                              prefix is accepted ("#80FF5722").
        align: one of "left" / "center" / "right" (horizontal only).
        border: True adds a thin black border to every cell in the
                range; False clears borders.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={range, cells_formatted, applied}` where `applied`
        is the list of attribute names that were actually changed.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "excel":
        return _write_error(kind, "office.set_cell_format requires a .xlsx file", resolved)

    applied: list[str] = []
    try:
        font_color_hex = _color_to_hex(font_color) if font_color else None
        bg_color_hex = _color_to_hex(bg_color) if bg_color else None
    except ValueError as exc:
        return _write_error(kind, str(exc), resolved)
    if align is not None and align not in {"left", "center", "right"}:
        return _write_error(kind, f"align must be left/center/right, got {align!r}", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        wb = load_workbook(target_path)
        try:
            ws = _resolve_sheet(wb, sheet)
            min_row, min_col, max_row, max_col = _parse_range(range)
            count = 0
            side = Side(border_style="thin", color="FF000000") if border else None
            border_style = (
                Border(left=side, right=side, top=side, bottom=side)
                if border
                else (Border() if border is False else None)
            )
            for r in range_iter(min_row, max_row):
                for c in range_iter(min_col, max_col):
                    cell_obj = ws.cell(row=r, column=c)
                    # Font — clone current font then patch fields the
                    # caller asked to change.
                    if any(v is not None for v in (bold, italic, font_size, font_color_hex)):
                        existing = cell_obj.font
                        cell_obj.font = Font(
                            name=existing.name,
                            size=font_size if font_size is not None else existing.size,
                            bold=bold if bold is not None else existing.bold,
                            italic=italic if italic is not None else existing.italic,
                            color=font_color_hex if font_color_hex is not None else existing.color,
                        )
                        for f in ("bold", "italic", "font_size", "font_color"):
                            if locals().get(f) is not None or (
                                f == "font_color" and font_color_hex
                            ):
                                if f not in applied:
                                    applied.append(f)
                    if bg_color_hex is not None:
                        cell_obj.fill = PatternFill(
                            fill_type="solid", start_color=bg_color_hex, end_color=bg_color_hex
                        )
                        if "bg_color" not in applied:
                            applied.append("bg_color")
                    if align is not None:
                        existing = cell_obj.alignment
                        cell_obj.alignment = Alignment(
                            horizontal=align,
                            vertical=existing.vertical,
                            wrap_text=existing.wrap_text,
                        )
                        if "align" not in applied:
                            applied.append("align")
                    if border_style is not None:
                        cell_obj.border = border_style
                        if "border" not in applied:
                            applied.append("border")
                    count += 1
            wb.save(tmp_path)
        finally:
            wb.close()
        return {"range": range, "cells_formatted": count, "applied": applied}

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


def range_iter(lo: int, hi: int):
    """Inclusive integer range — `range(lo, hi+1)` shorthand."""
    return range(lo, hi + 1)


@tool("office.merge_cells")
def office_merge_cells(
    path: str,
    sheet: str | None,
    range: str,
) -> dict[str, Any]:
    """Merge a rectangular range into a single cell.

    Writes to `<basename>.edited.xlsx`; original preserved.

    Args:
        path: .xlsx to edit.
        sheet: sheet name (None = active).
        range: A1-style range like "A1:C1" (header row span) or
               "A1:B4" (vertical+horizontal merge).

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={range}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "excel":
        return _write_error(kind, "office.merge_cells requires a .xlsx file", resolved)
    if ":" not in range:
        return _write_error(
            kind, "merge_cells requires a range like 'A1:C1', not a single cell", resolved
        )

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        wb = load_workbook(target_path)
        try:
            ws = _resolve_sheet(wb, sheet)
            ws.merge_cells(range)
            wb.save(tmp_path)
        finally:
            wb.close()
        return {"range": range}

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


@tool("office.add_row")
def office_add_row(
    path: str,
    sheet: str | None,
    values: list[Any],
    position: str | int = "append",
) -> dict[str, Any]:
    """Insert a row into a sheet.

    Writes to `<basename>.edited.xlsx`; original preserved.

    Args:
        path: .xlsx to edit.
        sheet: sheet name (None = active).
        values: list of cell values; the row is filled left-to-right
                from column A.
        position: "append" (default — adds after the last used row)
                  or an integer row number (1-indexed) to insert BEFORE
                  that row, shifting subsequent rows down.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={row, cells_written}`.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "excel":
        return _write_error(kind, "office.add_row requires a .xlsx file", resolved)
    if not isinstance(values, list):
        return _write_error(kind, "values must be a list", resolved)
    if position != "append" and not isinstance(position, int):
        return _write_error(kind, "position must be 'append' or an integer row number", resolved)

    target_path = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        wb = load_workbook(target_path)
        try:
            ws = _resolve_sheet(wb, sheet)
            if position == "append":
                ws.append(values)
                row_idx = ws.max_row
            else:
                row_idx = int(position)
                if row_idx < 1:
                    raise ValueError(f"row position must be >= 1, got {row_idx}")
                ws.insert_rows(row_idx)
                for i, v in enumerate(values, start=1):
                    ws.cell(row=row_idx, column=i, value=v)
            wb.save(tmp_path)
        finally:
            wb.close()
        return {"row": row_idx, "cells_written": len(values)}

    try:
        summary = _atomic_write(target_path, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target_path, kind, summary)


@tool("office.read_range")
def office_read_range(
    path: str,
    sheet: str | None,
    range: str,
) -> dict[str, Any]:
    """Read a specific cell range from a .xlsx as structured JSON.

    This is the precise complement to `office.read` (which dumps the
    whole workbook as markdown). Use this when you only need a subset
    — much smaller LLM context spend, and you get raw types (int /
    float / str / bool / None) plus the formula text for formula
    cells.

    Args:
        path: .xlsx to read (may be original or `.edited` copy).
        sheet: sheet name (None = active).
        range: A1-style range like "A1:C10" or single cell "B5".

    Returns:
        dict with:
          ok
          path, kind
          sheet (resolved sheet name)
          range (echoed back)
          values   — 2D list, computed values (formulas → cached result)
          formulas — 2D list, formula text where present, else None
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return {"ok": "false", "error": err}
    assert resolved is not None and kind is not None
    if kind != "excel":
        return {"ok": "false", "error": "office.read_range requires a .xlsx file", "path": resolved}
    try:
        # Two passes: data_only=True for cached values, data_only=False
        # for formula text. openpyxl can't give both in one open.
        wb_values = load_workbook(resolved, read_only=True, data_only=True)
        wb_formulas = load_workbook(resolved, read_only=True, data_only=False)
    except Exception as exc:
        return {
            "ok": "false",
            "path": resolved,
            "kind": kind,
            "error": f"failed to open .xlsx: {exc.__class__.__name__}: {exc}",
        }
    try:
        ws_values = _resolve_sheet(wb_values, sheet)
        ws_formulas = _resolve_sheet(wb_formulas, sheet)
        min_row, min_col, max_row, max_col = _parse_range(range)
        values_grid: list[list[Any]] = []
        formulas_grid: list[list[str | None]] = []
        for r in range_iter(min_row, max_row):
            v_row: list[Any] = []
            f_row: list[str | None] = []
            for c in range_iter(min_col, max_col):
                v_row.append(ws_values.cell(row=r, column=c).value)
                f_cell = ws_formulas.cell(row=r, column=c).value
                f_row.append(f_cell if isinstance(f_cell, str) and f_cell.startswith("=") else None)
            values_grid.append(v_row)
            formulas_grid.append(f_row)
    except Exception as exc:
        wb_values.close()
        wb_formulas.close()
        return {
            "ok": "false",
            "path": resolved,
            "kind": kind,
            "error": f"failed to read range {range!r}: {exc.__class__.__name__}: {exc}",
        }
    sheet_name = ws_values.title
    wb_values.close()
    wb_formulas.close()
    return {
        "ok": "true",
        "path": resolved,
        "kind": kind,
        "sheet": sheet_name,
        "range": range,
        "values": values_grid,
        "formulas": formulas_grid,
    }


# ═══════════════════════════════════════════════════════════════════════
# PPTX helpers
# ═══════════════════════════════════════════════════════════════════════


def _resolve_slide_layout(prs, layout_name: str | None):
    """Return a slide_layout by name; fall back to layout index 1
    ('Title and Content') when name is None or not found."""
    if layout_name:
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
    # python-pptx default deck has these layouts in this order:
    #   0 Title Slide   1 Title and Content   2 Section Header
    #   3 Two Content   4 Comparison          5 Title Only
    #   6 Blank         7 Content with Caption 8 Picture with Caption
    return prs.slide_layouts[1]


def _set_placeholder_text(placeholder, value: str) -> None:
    """Replace a placeholder's text frame with `value`, preserving the
    first paragraph's formatting. python-pptx's `placeholder.text = ...`
    setter does this but drops bullet formatting; we want a single
    paragraph rewrite that survives the round-trip."""
    placeholder.text_frame.text = value


def _set_placeholder_bullets(placeholder, bullets: list[str]) -> None:
    """Replace a placeholder's text frame with one paragraph per bullet."""
    tf = placeholder.text_frame
    tf.clear()  # leaves one empty paragraph
    if not bullets:
        return
    # The cleared text_frame has exactly one paragraph; reuse it for
    # the first bullet, then add the rest.
    tf.paragraphs[0].text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet


def _find_content_placeholder(slide):
    """Return the body placeholder (anything not the title). None when
    the layout has no body slot — caller should handle by adding a
    text box or returning an error."""
    title_shape = slide.shapes.title
    title_elem = title_shape._element if title_shape is not None else None
    for shape in slide.placeholders:
        if title_elem is not None and shape._element is title_elem:
            continue
        return shape
    return None


def _set_slide_content(slide, title: str | None, bullets: list[str] | None) -> None:
    """Apply title/bullets to a slide, leaving unset fields alone."""
    if title is not None:
        if slide.shapes.title is None:
            raise ValueError("slide layout has no title placeholder")
        _set_placeholder_text(slide.shapes.title, title)
    if bullets is not None:
        body = _find_content_placeholder(slide)
        if body is None:
            raise ValueError("slide layout has no body placeholder for bullets")
        _set_placeholder_bullets(body, bullets)


def _validate_slide_index(prs, index: int) -> None:
    n = len(prs.slides)
    if not (0 <= index < n):
        raise ValueError(f"slide index {index} out of range (deck has {n} slides)")


# ═══════════════════════════════════════════════════════════════════════
# PPTX write tools
# ═══════════════════════════════════════════════════════════════════════


@tool("office.create_pptx")
def office_create_pptx(path: str, title: str | None = None) -> dict[str, Any]:
    """Create a new blank PowerPoint deck at `path` with one title slide.

    Special-case among write tools: this writes DIRECTLY to `path`
    (not to a `.edited.pptx` copy) — there's no pre-existing original
    to protect when you're creating from scratch. Subsequent edits
    via office.add_slide / office.update_slide / etc. follow the
    standard copy-on-first-write pattern and land in
    `<basename>.edited.pptx`.

    Args:
        path: absolute path where the new .pptx should land (extension
              must be .pptx).
        title: optional title for the first slide. When omitted the
               slide is created with an empty title placeholder.

    Returns:
        dict with `ok`, `original_path`, `edited_path` (= path,
        because we wrote the original), `kind="powerpoint"`,
        `summary={slide_count: 1}`.
    """
    if not path:
        return _write_error(None, "path required")
    resolved = os.path.abspath(os.path.expanduser(path))
    if Path(resolved).suffix.lower() != ".pptx":
        return _write_error(None, "office.create_pptx requires a .pptx path")
    if os.path.exists(resolved):
        return _write_error(
            "powerpoint",
            f"file already exists: {resolved} (use office.update_slide to edit it)",
            resolved,
        )
    try:
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        if title is not None and slide.shapes.title is not None:
            _set_placeholder_text(slide.shapes.title, title)
        # Ensure the parent directory exists (project-relative paths).
        Path(resolved).parent.mkdir(parents=True, exist_ok=True)
        prs.save(resolved)
        _verify_file(resolved, "powerpoint")
    except Exception as exc:
        # Clean up partial write so a half-created file doesn't linger.
        if os.path.exists(resolved):
            try:
                os.remove(resolved)
            except OSError:
                pass
        return _write_error(
            "powerpoint", f"create failed: {exc.__class__.__name__}: {exc}", resolved
        )
    return {
        "ok": "true",
        "original_path": resolved,
        "edited_path": resolved,  # the create IS the working file
        "kind": "powerpoint",
        "summary": {"slide_count": 1, "title": title or ""},
    }


@tool("office.add_slide")
def office_add_slide(
    path: str,
    layout: str | None = None,
    title: str | None = None,
    bullets: list[str] | None = None,
) -> dict[str, Any]:
    """Append a new slide to a .pptx. Writes to `<basename>.edited.pptx`.

    Args:
        path: .pptx to edit.
        layout: python-pptx layout name (e.g. "Title and Content",
                "Title Slide", "Section Header", "Two Content", "Blank").
                Default: "Title and Content".
        title: optional title text.
        bullets: optional list of strings; each becomes one paragraph
                 (bullet) in the body placeholder.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={index, layout}` where `index` is the new slide's
        0-based position.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.add_slide requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        slide_layout = _resolve_slide_layout(prs, layout)
        slide = prs.slides.add_slide(slide_layout)
        _set_slide_content(slide, title, bullets)
        prs.save(tmp_path)
        return {"index": len(prs.slides) - 1, "layout": slide_layout.name}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.update_slide")
def office_update_slide(
    path: str,
    index: int,
    title: str | None = None,
    bullets: list[str] | None = None,
    layout: str | None = None,
) -> dict[str, Any]:
    """Update an existing slide's title / bullets / layout.

    Writes to `<basename>.edited.pptx`. Any argument left as None is
    untouched. NOTE: changing `layout` re-creates the slide on the new
    layout, so any custom shapes you added previously are lost — only
    the title + body bullets carry over.

    Args:
        path: .pptx to edit.
        index: 0-based slide index.
        title, bullets, layout: see add_slide.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.update_slide requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        if layout is not None:
            # Layout change: remove the existing slide, insert a new one
            # at the same position with the new layout.
            existing = prs.slides[index]
            existing_title = _slide_title_text(existing)
            existing_bullets = _slide_bullets(existing)
            new_title = title if title is not None else existing_title
            new_bullets = bullets if bullets is not None else existing_bullets
            new_layout = _resolve_slide_layout(prs, layout)
            new_slide = prs.slides.add_slide(new_layout)
            _set_slide_content(new_slide, new_title, new_bullets)
            # Move the new slide (currently last) to the original index,
            # then delete the original.
            _move_slide(prs, len(prs.slides) - 1, index + 1)  # insert after the original
            _delete_slide(prs, index)
            applied_layout = new_layout.name
        else:
            slide = prs.slides[index]
            _set_slide_content(slide, title, bullets)
            applied_layout = slide.slide_layout.name if slide.slide_layout else ""
        prs.save(tmp_path)
        return {"index": index, "layout": applied_layout}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


def _delete_slide(prs, index: int) -> None:
    """python-pptx has no public delete; manipulate sldIdLst directly."""
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    to_remove = slides[index]
    # Drop the relationship so PowerPoint doesn't complain about a
    # dangling rId on open.
    rels_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    rid = to_remove.attrib[rels_attr]
    prs.part.drop_rel(rid)
    sld_id_lst.remove(to_remove)


def _move_slide(prs, from_index: int, to_index: int) -> None:
    """Reorder slides via sldIdLst manipulation. to_index is the
    desired final position; negative indices not supported."""
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    if not (0 <= from_index < len(slides)):
        raise ValueError(f"from_index {from_index} out of range")
    if not (0 <= to_index < len(slides)):
        raise ValueError(f"to_index {to_index} out of range")
    moving = slides[from_index]
    sld_id_lst.remove(moving)
    sld_id_lst.insert(to_index, moving)


@tool("office.delete_slide")
def office_delete_slide(path: str, index: int) -> dict[str, Any]:
    """Delete the slide at `index` (0-based). Writes to `.edited.pptx`."""
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.delete_slide requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        deleted_title = _slide_title_text(prs.slides[index])
        _delete_slide(prs, index)
        prs.save(tmp_path)
        return {
            "deleted_index": index,
            "deleted_title": deleted_title,
            "slide_count": len(prs.slides),
        }

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.reorder_slides")
def office_reorder_slides(path: str, from_index: int, to_index: int) -> dict[str, Any]:
    """Move slide from `from_index` to `to_index` (both 0-based).
    Writes to `.edited.pptx`."""
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.reorder_slides requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _move_slide(prs, from_index, to_index)
        prs.save(tmp_path)
        return {"from_index": from_index, "to_index": to_index}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.set_slide_title")
def office_set_slide_title(path: str, index: int, title: str) -> dict[str, Any]:
    """Set just the title text of slide `index`. Writes to `.edited.pptx`."""
    return office_update_slide.func(path=path, index=index, title=title)


@tool("office.set_slide_bullets")
def office_set_slide_bullets(path: str, index: int, bullets: list[str]) -> dict[str, Any]:
    """Replace the body bullets of slide `index`. Writes to `.edited.pptx`.

    `bullets` is a list of strings — each becomes one paragraph in the
    body placeholder. Pass `[]` to clear the body."""
    return office_update_slide.func(path=path, index=index, bullets=bullets)


@tool("office.set_slide_notes")
def office_set_slide_notes(path: str, index: int, notes: str) -> dict[str, Any]:
    """Set the presenter notes for slide `index`. Writes to `.edited.pptx`."""
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.set_slide_notes requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        slide = prs.slides[index]
        # has_notes_slide is True if notes have ever been set; accessing
        # notes_slide on a slide-without-notes auto-creates the notes slide.
        slide.notes_slide.notes_text_frame.text = notes
        prs.save(tmp_path)
        return {"index": index, "notes_excerpt": notes[:60]}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.add_image_to_slide")
def office_add_image_to_slide(
    path: str,
    index: int,
    image_path: str,
    left_in: float | None = None,
    top_in: float | None = None,
    width_in: float | None = None,
) -> dict[str, Any]:
    """Insert an image into slide `index`. Writes to `.edited.pptx`.

    Args:
        path: .pptx to edit.
        index: 0-based slide index.
        image_path: absolute path to an image file (.png, .jpg, .gif).
                    Must already exist on disk.
        left_in, top_in, width_in: position + size in inches. When
                                   omitted the image is centered on
                                   the slide at ~50% width.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.add_image_to_slide requires a .pptx file", resolved)
    img_resolved = os.path.abspath(os.path.expanduser(image_path))
    if not os.path.isfile(img_resolved):
        return _write_error(kind, f"image file not found: {img_resolved}", resolved)

    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        slide = prs.slides[index]
        # Slide width is usually 10in; height ~7.5in. Default: 5in wide,
        # centered horizontally, ~30% from top.
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        w = Inches(width_in) if width_in is not None else Inches(5)
        left = Inches(left_in) if left_in is not None else (slide_w - w) // 2
        top = Inches(top_in) if top_in is not None else slide_h // 4
        pic = slide.shapes.add_picture(img_resolved, left, top, width=w)
        prs.save(tmp_path)
        return {
            "index": index,
            "image": img_resolved,
            "width_in": float(pic.width) / 914400,  # EMU → inches
        }

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


# ═══════════════════════════════════════════════════════════════════════
# PPTX read tool
# ═══════════════════════════════════════════════════════════════════════


@tool("office.read_slides")
def office_read_slides(path: str) -> dict[str, Any]:
    """Return per-slide structured data (title, bullets, notes, …).

    Use this instead of `office.read` when you want LLM-friendly slide
    semantics (one entry per slide with title + bullets list) rather
    than the markitdown markdown dump. The frontend preview panel
    also calls this through a small HTTP endpoint to render the
    outline view.

    Args:
        path: .pptx to read (may be original or `.edited` copy).

    Returns:
        dict with:
          ok, path, kind="powerpoint"
          slides (list of {index, layout, title, bullets, notes,
                  shape_count, image_count})
          slide_count
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return {"ok": "false", "error": err}
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return {
            "ok": "false",
            "path": resolved,
            "error": "office.read_slides requires a .pptx file",
        }
    try:
        details = _outline_pptx(resolved)
    except Exception as exc:
        return {
            "ok": "false",
            "path": resolved,
            "kind": kind,
            "error": f"failed to read .pptx: {exc.__class__.__name__}: {exc}",
        }
    return {"ok": "true", "path": resolved, "kind": kind, **details}


# Phase 1 surface (read-only).
OFFICE_READ_TOOLS = [office_read, office_outline, office_read_range, office_read_slides]

# Phase 2 + Phase 3 surface (writes, all copy-on-first-write — except
# office.create_pptx which produces the original file).
OFFICE_WRITE_TOOLS = [
    office_find_replace,
    office_insert_paragraph,
    office_update_paragraph,
    office_delete_paragraph,
    office_apply_style,
    office_set_cells,
    office_set_formula,
    office_set_cell_format,
    office_merge_cells,
    office_add_row,
    office_create_pptx,
    office_add_slide,
    office_update_slide,
    office_delete_slide,
    office_reorder_slides,
    office_set_slide_title,
    office_set_slide_bullets,
    office_set_slide_notes,
    office_add_image_to_slide,
]


# Re-exported for tests that want to construct paths the same way the
# tools do without poking at the private helper.
def edited_copy_path(original_path: str) -> str:
    return _resolve_write_target(original_path)


# Silence unused-import warnings for symbols re-exported by tests.
__all__ = [
    "OFFICE_READ_TOOLS",
    "OFFICE_WRITE_TOOLS",
    "edited_copy_path",
    "get_column_letter",  # tests construct range refs
    "office_add_image_to_slide",
    "office_add_row",
    "office_add_slide",
    "office_apply_style",
    "office_create_pptx",
    "office_delete_paragraph",
    "office_delete_slide",
    "office_find_replace",
    "office_insert_paragraph",
    "office_merge_cells",
    "office_outline",
    "office_read",
    "office_read_range",
    "office_read_slides",
    "office_reorder_slides",
    "office_set_cell_format",
    "office_set_cells",
    "office_set_formula",
    "office_set_slide_bullets",
    "office_set_slide_notes",
    "office_set_slide_title",
    "office_update_paragraph",
    "office_update_slide",
]
