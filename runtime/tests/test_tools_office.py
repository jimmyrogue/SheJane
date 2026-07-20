"""Unit tests for office.read and office.outline.

The fixtures here are built on-the-fly with python-docx / openpyxl so we
don't have to check binary files into the repo. Each test owns its own
tmp_path-scoped file.

Office tools don't proxy through any external service (markitdown +
python-docx + openpyxl all run in-process), so these are vanilla unit
tests — no MockTransport, no settings override needed.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook

from shejane_runtime.tools.office import (
    OFFICE_READ_TOOLS,
    OFFICE_WRITE_TOOLS,
    edited_copy_path,
    office_add_image_to_slide,
    office_add_row,
    office_add_slide,
    office_apply_style,
    office_create_pptx,
    office_delete_paragraph,
    office_delete_slide,
    office_find_replace,
    office_insert_paragraph,
    office_merge_cells,
    office_outline,
    office_read,
    office_read_range,
    office_read_slides,
    office_reorder_slides,
    office_set_cell_format,
    office_set_cells,
    office_set_formula,
    office_set_slide_bullets,
    office_set_slide_notes,
    office_set_slide_title,
    office_update_paragraph,
    office_update_slide,
)


def _make_docx(tmp_path: Path) -> Path:
    """Build a tiny .docx with heading + paragraphs + a table."""
    p = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph("First paragraph with some prose.")
    doc.add_heading("Methodology", level=2)
    doc.add_paragraph("Methodology details here.")
    doc.add_paragraph("Another methodology paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Year"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "2025"
    table.rows[1].cells[1].text = "42"
    doc.save(p)
    return p


def _make_xlsx(tmp_path: Path) -> Path:
    """Build a tiny .xlsx with two sheets."""
    p = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sales"
    ws1.append(["Quarter", "Revenue"])
    ws1.append(["Q1", 100])
    ws1.append(["Q2", 150])
    ws2 = wb.create_sheet("Costs")
    ws2.append(["Item", "Amount"])
    ws2.append(["Rent", 50])
    wb.save(p)
    return p


def test_office_read_docx_returns_markdown(tmp_path: Path) -> None:
    """office.read on a .docx → markdown with the headings preserved."""
    path = _make_docx(tmp_path)
    result = office_read.func(path=str(path))
    assert result["ok"] == "true"
    assert result["kind"] == "word"
    md = result["markdown"]
    # Heading text should survive in the markdown output. markitdown's
    # exact heading-level mapping varies; assert the text is there.
    assert "Executive Summary" in md
    assert "Methodology" in md
    # Table contents should also surface.
    assert "Year" in md and "Value" in md
    # Not truncated for a tiny doc.
    assert result["truncated"] == "false"


def test_office_read_xlsx_returns_markdown(tmp_path: Path) -> None:
    """office.read on a .xlsx → markdown with both sheets."""
    path = _make_xlsx(tmp_path)
    result = office_read.func(path=str(path))
    assert result["ok"] == "true"
    assert result["kind"] == "excel"
    md = result["markdown"]
    # Both sheet names should appear (markitdown sections by sheet).
    assert "Sales" in md
    assert "Costs" in md
    # Cell contents should surface.
    assert "Q1" in md and "100" in md


def test_office_read_missing_file_returns_error(tmp_path: Path) -> None:
    """office.read on a non-existent path → ok=false, no crash."""
    result = office_read.func(path=str(tmp_path / "missing.docx"))
    assert result["ok"] == "false"
    assert "not found" in result["error"]


def test_office_read_wrong_extension_returns_error(tmp_path: Path) -> None:
    """office.read on a .txt → ok=false steering the LLM to read_file."""
    txt = tmp_path / "note.txt"
    txt.write_text("plain text")
    result = office_read.func(path=str(txt))
    assert result["ok"] == "false"
    # Error message should mention the supported extensions so the LLM
    # can self-correct on the next call.
    assert ".docx" in result["error"] and ".xlsx" in result["error"]


def test_office_read_empty_path_returns_error() -> None:
    """office.read with empty path returns a structured error."""
    result = office_read.func(path="")
    assert result["ok"] == "false"
    assert result["error"] == "path required"


def test_office_read_rejects_paths_outside_run_workspace(tmp_path: Path) -> None:
    workspace = tmp_path / "workspace"
    outside = tmp_path / "outside"
    workspace.mkdir()
    outside.mkdir()
    outside_doc = _make_docx(outside)

    result = office_read.invoke(
        {"path": str(outside_doc)},
        config={"configurable": {"workspace_root": str(workspace)}},
    )

    assert result["ok"] == "false"
    assert "outside workspace" in result["error"]


def test_office_read_rejects_paths_when_run_has_no_workspace(tmp_path: Path) -> None:
    path = _make_docx(tmp_path)

    result = office_read.invoke(
        {"path": str(path)},
        config={"configurable": {"workspace_root": ""}},
    )

    assert result["ok"] == "false"
    assert result["error"] == "no workspace open"


def test_office_outline_docx_lists_headings(tmp_path: Path) -> None:
    """office.outline on .docx surfaces heading text + paragraph/table counts."""
    path = _make_docx(tmp_path)
    result = office_outline.func(path=str(path))
    assert result["ok"] == "true"
    assert result["kind"] == "word"
    heading_texts = [h["text"] for h in result["headings"]]
    assert "Executive Summary" in heading_texts
    assert "Methodology" in heading_texts
    # Two top-level + one subsection means at least one heading per level.
    levels = {h["level"] for h in result["headings"]}
    assert 1 in levels and 2 in levels
    # Paragraphs: 2 headings + 3 body paragraphs (heading paragraphs count
    # too in python-docx's iteration); >= 5 is the safe bound.
    assert result["paragraph_count"] >= 5
    assert result["table_count"] == 1


def test_office_outline_xlsx_lists_sheets(tmp_path: Path) -> None:
    """office.outline on .xlsx surfaces sheet names + dimensions."""
    path = _make_xlsx(tmp_path)
    result = office_outline.func(path=str(path))
    assert result["ok"] == "true"
    assert result["kind"] == "excel"
    sheet_names = [s["name"] for s in result["sheets"]]
    assert sheet_names == ["Sales", "Costs"]
    sales = next(s for s in result["sheets"] if s["name"] == "Sales")
    assert sales["rows"] == 3  # header + 2 data rows
    assert sales["columns"] == 2


def test_office_outline_wrong_extension_returns_error(tmp_path: Path) -> None:
    """office.outline rejects non-office files just like office.read."""
    txt = tmp_path / "x.csv"
    txt.write_text("a,b\n1,2\n")
    result = office_outline.func(path=str(txt))
    assert result["ok"] == "false"
    assert "unsupported extension" in result["error"]


def test_office_tools_registered_under_canonical_names() -> None:
    """Sanity: OFFICE_READ_TOOLS exposes the read tools by their dotted names."""
    names = sorted(t.name for t in OFFICE_READ_TOOLS)
    assert names == [
        "office.outline",
        "office.read",
        "office.read_range",
        "office.read_slides",
    ]


def test_office_read_truncates_large_output(tmp_path: Path) -> None:
    """office.read above the 60_000 char cap → truncated="true" + the cap."""
    from shejane_runtime.tools import office as office_module

    # Build a .docx with way more than 60k chars of body text by repeating
    # a paragraph many times. Cheaper than the LLM cap test using fixtures.
    p = tmp_path / "big.docx"
    doc = Document()
    para = "Lorem ipsum dolor sit amet, consectetur adipiscing elit. " * 50  # ~2900 chars/para
    for _ in range(25):  # ~72500 chars total before markdown overhead
        doc.add_paragraph(para)
    doc.save(p)
    result = office_read.func(path=str(p))
    assert result["ok"] == "true"
    assert result["truncated"] == "true"
    # The cap is enforced on the returned markdown; allow some slack for
    # the trailing "…(truncated, full size = N chars)" suffix.
    assert len(result["markdown"]) >= office_module._MARKDOWN_CHAR_CAP
    assert len(result["markdown"]) < office_module._MARKDOWN_CHAR_CAP + 200


@pytest.mark.parametrize("ext", [".doc", ".xls", ".pdf"])
def test_office_read_rejects_non_supported_office_extensions(tmp_path: Path, ext: str) -> None:
    """office.read rejects .doc / .xls / .pdf — we only support OOXML
    formats (.docx / .xlsx / .pptx as of Phase 3). The legacy binary
    formats would need a separate parser."""
    p = tmp_path / f"x{ext}"
    p.write_bytes(b"\x00")  # dummy bytes, validation happens on extension first
    result = office_read.func(path=str(p))
    assert result["ok"] == "false"
    assert ext in result["error"] or "unsupported extension" in result["error"]


# ═══════════════════════════════════════════════════════════════════════
# Phase 2: write tools (copy-on-first-write)
# ═══════════════════════════════════════════════════════════════════════


def _sha256(path: Path) -> str:
    import hashlib

    return hashlib.sha256(path.read_bytes()).hexdigest()


def _docx_paragraphs(path: Path) -> list[str]:
    return [p.text for p in Document(path).paragraphs]


def _xlsx_cell(path: Path, sheet: str, cell: str):
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=False)
    try:
        return wb[sheet][cell].value
    finally:
        wb.close()


def test_write_tools_registered_under_canonical_names() -> None:
    """Sanity: OFFICE_WRITE_TOOLS exposes all 19 write tools by name
    (10 Phase 2 docx/xlsx + 9 Phase 3 pptx)."""
    names = sorted(t.name for t in OFFICE_WRITE_TOOLS)
    assert names == sorted(
        [
            # Phase 2 — docx/xlsx
            "office.find_replace",
            "office.insert_paragraph",
            "office.update_paragraph",
            "office.delete_paragraph",
            "office.apply_style",
            "office.set_cells",
            "office.set_formula",
            "office.set_cell_format",
            "office.merge_cells",
            "office.add_row",
            # Phase 3 — pptx
            "office.create_pptx",
            "office.add_slide",
            "office.update_slide",
            "office.delete_slide",
            "office.reorder_slides",
            "office.set_slide_title",
            "office.set_slide_bullets",
            "office.set_slide_notes",
            "office.add_image_to_slide",
        ]
    )


def test_edited_copy_path_helper() -> None:
    """`<basename>.edited.<ext>` lives in the same dir; already-edited
    paths are returned unchanged."""
    assert edited_copy_path("/x/y/report.docx") == "/x/y/report.edited.docx"
    assert edited_copy_path("/x/y/report.edited.docx") == "/x/y/report.edited.docx"
    assert edited_copy_path("/x/y/q4.xlsx") == "/x/y/q4.edited.xlsx"


def test_find_replace_writes_to_copy_and_preserves_original(tmp_path: Path) -> None:
    """Core contract: original SHA256 unchanged; .edited contains the replacement."""
    p = _make_docx(tmp_path)
    sha_before = _sha256(p)
    result = office_find_replace.func(
        path=str(p),
        find="First paragraph",
        replace="Updated first sentence",
    )
    assert result["ok"] == "true", result
    edited = Path(result["edited_path"])
    assert edited.exists()
    assert edited != p
    # Original byte-for-byte unchanged.
    assert _sha256(p) == sha_before
    # Edited copy has the new text.
    paragraphs = _docx_paragraphs(edited)
    assert any("Updated first sentence" in para for para in paragraphs)
    assert all("First paragraph with some prose." not in para for para in paragraphs)
    assert result["summary"]["replaced"] == 1


def test_find_replace_respects_count_limit(tmp_path: Path) -> None:
    """count=N caps the global number of replacements."""
    p = tmp_path / "repeats.docx"
    doc = Document()
    for _ in range(5):
        doc.add_paragraph("apple banana apple")
    doc.save(p)
    result = office_find_replace.func(path=str(p), find="apple", replace="X", count=3)
    assert result["ok"] == "true"
    assert result["summary"]["replaced"] == 3
    text = "\n".join(_docx_paragraphs(Path(result["edited_path"])))
    # Started with 5 paragraphs × 2 = 10 "apple"; 3 replaced → 7 remain.
    assert text.count("apple") == 7
    assert text.count("X") == 3


def test_repeat_write_targets_same_edited_file(tmp_path: Path) -> None:
    """Two writes against the original land in ONE `.edited.docx`, the
    second building on the first."""
    p = _make_docx(tmp_path)
    r1 = office_find_replace.func(path=str(p), find="Methodology", replace="Approach")
    assert r1["ok"] == "true"
    edited = Path(r1["edited_path"])
    sha_after_first = _sha256(edited)

    r2 = office_find_replace.func(path=str(p), find="Executive", replace="Top-line")
    assert r2["ok"] == "true"
    # Same edited path both times.
    assert r2["edited_path"] == r1["edited_path"]
    # Edited file changed between writes (we built on top of the first).
    assert _sha256(edited) != sha_after_first
    # Second edit incorporates BOTH changes.
    texts = "\n".join(_docx_paragraphs(edited))
    assert "Approach" in texts
    assert "Top-line" in texts
    assert "Methodology" not in texts
    assert "Executive" not in texts


def test_insert_paragraph_after_anchor(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    result = office_insert_paragraph.func(
        path=str(p),
        anchor="Methodology",
        content="Inserted sentence.",
        position="after",
        style=None,
    )
    assert result["ok"] == "true", result
    edited_paragraphs = _docx_paragraphs(Path(result["edited_path"]))
    # The inserted paragraph appears AFTER the Methodology heading.
    idx_methodology = next(i for i, t in enumerate(edited_paragraphs) if "Methodology" in t)
    assert edited_paragraphs[idx_methodology + 1] == "Inserted sentence."


def test_insert_paragraph_at_end_with_style(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    before_count = len(_docx_paragraphs(p))
    result = office_insert_paragraph.func(
        path=str(p),
        anchor="__END__",
        content="Conclusion",
        style="Heading 1",
    )
    assert result["ok"] == "true"
    paragraphs = _docx_paragraphs(Path(result["edited_path"]))
    assert paragraphs[-1] == "Conclusion"
    assert len(paragraphs) == before_count + 1


def test_insert_paragraph_missing_anchor_returns_error(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    result = office_insert_paragraph.func(
        path=str(p),
        anchor="this-string-doesnt-exist",
        content="…",
    )
    assert result["ok"] == "false"
    assert "anchor not found" in result["error"]


def test_update_paragraph_replaces_full_text(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    result = office_update_paragraph.func(
        path=str(p),
        target="First paragraph",
        content="Entirely rewritten paragraph body.",
    )
    assert result["ok"] == "true"
    paragraphs = _docx_paragraphs(Path(result["edited_path"]))
    assert any("Entirely rewritten paragraph body." in t for t in paragraphs)
    assert all("First paragraph with some prose." not in t for t in paragraphs)


def test_delete_paragraph_removes_one(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    before = _docx_paragraphs(p)
    result = office_delete_paragraph.func(path=str(p), target="Methodology details")
    assert result["ok"] == "true"
    after = _docx_paragraphs(Path(result["edited_path"]))
    assert len(after) == len(before) - 1
    assert not any("Methodology details" in t for t in after)


def test_apply_style_changes_style(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    result = office_apply_style.func(
        path=str(p),
        target="First paragraph",
        style="Heading 2",
    )
    assert result["ok"] == "true"
    assert result["summary"]["new_style"] == "Heading 2"
    # Reopen and verify the style stuck.
    doc = Document(Path(result["edited_path"]))
    match = next(p for p in doc.paragraphs if "First paragraph" in p.text)
    assert match.style.name == "Heading 2"


def test_set_cells_writes_2d_block(tmp_path: Path) -> None:
    p = _make_xlsx(tmp_path)
    sha_before = _sha256(p)
    result = office_set_cells.func(
        path=str(p),
        sheet="Sales",
        range="C1:D2",
        values=[["X", "Y"], [10, 20]],
    )
    assert result["ok"] == "true", result
    # Original untouched.
    assert _sha256(p) == sha_before
    edited = Path(result["edited_path"])
    assert _xlsx_cell(edited, "Sales", "C1") == "X"
    assert _xlsx_cell(edited, "Sales", "D2") == 20
    # Unrelated cells unaffected.
    assert _xlsx_cell(edited, "Sales", "A1") == "Quarter"
    assert result["summary"]["cells_written"] == 4


def test_set_formula_writes_formula_literal(tmp_path: Path) -> None:
    p = _make_xlsx(tmp_path)
    result = office_set_formula.func(
        path=str(p),
        sheet="Sales",
        cell="C5",
        formula="=SUM(B2:B3)",
    )
    assert result["ok"] == "true"
    # data_only=False reads the formula text (openpyxl can't evaluate;
    # the cached value is None until Excel opens the file).
    assert _xlsx_cell(Path(result["edited_path"]), "Sales", "C5") == "=SUM(B2:B3)"


def test_set_formula_rejects_missing_leading_equals(tmp_path: Path) -> None:
    p = _make_xlsx(tmp_path)
    result = office_set_formula.func(path=str(p), sheet="Sales", cell="C1", formula="SUM(A1:A3)")
    assert result["ok"] == "false"
    assert "must start with '='" in result["error"]


def test_set_cell_format_applies_bold_and_color(tmp_path: Path) -> None:
    from openpyxl import load_workbook

    p = _make_xlsx(tmp_path)
    result = office_set_cell_format.func(
        path=str(p),
        sheet="Sales",
        range="A1:B1",
        bold=True,
        font_color="#FF5722",
        bg_color="#FFFFCC",
    )
    assert result["ok"] == "true"
    wb = load_workbook(Path(result["edited_path"]))
    try:
        ws = wb["Sales"]
        cell = ws["A1"]
        assert cell.font.bold is True
        assert "FF5722" in (cell.font.color.rgb or "").upper()
        # PatternFill stores color object; coerce to string for compare.
        fill_color = str(cell.fill.start_color.rgb or "").upper()
        assert "FFFFCC" in fill_color
    finally:
        wb.close()
    assert "bold" in result["summary"]["applied"]
    assert result["summary"]["cells_formatted"] == 2


def test_set_cell_format_rejects_invalid_color(tmp_path: Path) -> None:
    p = _make_xlsx(tmp_path)
    result = office_set_cell_format.func(
        path=str(p), sheet="Sales", range="A1", font_color="not-a-hex"
    )
    assert result["ok"] == "false"
    assert "hex" in result["error"]


def test_merge_cells_records_merge(tmp_path: Path) -> None:
    from openpyxl import load_workbook

    p = _make_xlsx(tmp_path)
    result = office_merge_cells.func(path=str(p), sheet="Sales", range="A1:B1")
    assert result["ok"] == "true"
    wb = load_workbook(Path(result["edited_path"]))
    try:
        ranges = [str(r) for r in wb["Sales"].merged_cells.ranges]
        assert "A1:B1" in ranges
    finally:
        wb.close()


def test_add_row_append_and_insert(tmp_path: Path) -> None:
    p = _make_xlsx(tmp_path)

    # Append.
    r_append = office_add_row.func(
        path=str(p), sheet="Sales", values=["Q3", 200], position="append"
    )
    assert r_append["ok"] == "true"
    edited = Path(r_append["edited_path"])
    assert _xlsx_cell(edited, "Sales", "A4") == "Q3"
    assert _xlsx_cell(edited, "Sales", "B4") == 200

    # Insert at row 2 (before existing Q1 data).
    r_insert = office_add_row.func(path=str(p), sheet="Sales", values=["Q0", 50], position=2)
    assert r_insert["ok"] == "true"
    # row 2 is now Q0, row 3 (previously Q1) shifted down
    edited = Path(r_insert["edited_path"])
    assert _xlsx_cell(edited, "Sales", "A2") == "Q0"
    assert _xlsx_cell(edited, "Sales", "A3") == "Q1"


def test_read_range_returns_values_and_formulas(tmp_path: Path) -> None:
    p = _make_xlsx(tmp_path)
    # Plant a formula first so we have something to read.
    office_set_formula.func(path=str(p), sheet="Sales", cell="C2", formula="=B2*2")
    edited = Path(p).with_name("sample.edited.xlsx")
    assert edited.exists()

    result = office_read_range.func(path=str(edited), sheet="Sales", range="A1:C3")
    assert result["ok"] == "true"
    assert result["values"][0] == ["Quarter", "Revenue", None]
    assert result["values"][1][0] == "Q1"
    assert result["values"][1][1] == 100
    # Formula text surfaced at C2.
    assert result["formulas"][1][2] == "=B2*2"
    # Non-formula cells return None in the formulas grid.
    assert result["formulas"][0][0] is None


def test_read_range_rejects_docx(tmp_path: Path) -> None:
    p = _make_docx(tmp_path)
    result = office_read_range.func(path=str(p), sheet=None, range="A1:C3")
    assert result["ok"] == "false"
    assert ".xlsx" in result["error"]


def test_write_tool_rejects_unsupported_kind(tmp_path: Path) -> None:
    """find_replace is docx-only; calling on a .xlsx returns ok=false
    BEFORE creating any copy."""
    p = _make_xlsx(tmp_path)
    result = office_find_replace.func(path=str(p), find="x", replace="y")
    assert result["ok"] == "false"
    assert "requires a .docx" in result["error"]
    # Crucially: no .edited.xlsx file was created.
    assert not (tmp_path / "sample.edited.xlsx").exists()


def test_atomic_write_keeps_target_intact_on_verification_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """If verification fails after the tmp write, the target keeps its
    previous (last-successful) content and no tmp file is left behind."""
    from shejane_runtime.tools import office as office_module

    p = _make_docx(tmp_path)
    # First do a real edit so .edited.docx exists with known content.
    r1 = office_find_replace.func(path=str(p), find="Methodology", replace="Approach")
    assert r1["ok"] == "true"
    target = Path(r1["edited_path"])
    sha_known_good = _sha256(target)

    # Now force _verify_file to raise the next time it runs — simulates
    # a corrupted-on-save scenario.
    def _boom(path: str, kind: str) -> None:
        raise RuntimeError("simulated corruption")

    monkeypatch.setattr(office_module, "_verify_file", _boom)

    r2 = office_find_replace.func(path=str(p), find="Executive", replace="Top-line")
    assert r2["ok"] == "false"
    # Target file is unchanged (still the last known good).
    assert _sha256(target) == sha_known_good
    # No .tmp file lingering.
    assert not (tmp_path / "sample.edited.docx.tmp").exists()


# ═══════════════════════════════════════════════════════════════════════
# Phase 3: PPT tools
# ═══════════════════════════════════════════════════════════════════════


def _pptx_slide_titles(path: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(path)
    out: list[str] = []
    for slide in prs.slides:
        title = slide.shapes.title
        out.append(title.text_frame.text.strip() if title is not None else "")
    return out


def _pptx_slide_count(path: Path) -> int:
    from pptx import Presentation

    return len(Presentation(path).slides)


def test_create_pptx_writes_original_directly_no_edited(tmp_path: Path) -> None:
    """create_pptx is the special-case: writes to `<path>` (not
    `<path>.edited.pptx`) because there's no original to protect."""
    p = tmp_path / "deck.pptx"
    result = office_create_pptx.func(path=str(p), title="Hello")
    assert result["ok"] == "true", result
    assert result["edited_path"] == str(p)
    assert p.exists()
    # Should NOT have created a .edited copy yet.
    assert not (tmp_path / "deck.edited.pptx").exists()
    assert _pptx_slide_count(p) == 1
    assert _pptx_slide_titles(p) == ["Hello"]


def test_create_pptx_refuses_existing_file(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="A")
    second = office_create_pptx.func(path=str(p), title="B")
    assert second["ok"] == "false"
    assert "already exists" in second["error"]


def test_add_slide_creates_edited_copy_and_preserves_original(tmp_path: Path) -> None:
    """First add_slide triggers copy-on-first-write — original is frozen."""
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="Cover")
    sha_orig = _sha256(p)
    result = office_add_slide.func(path=str(p), title="Slide 2", bullets=["bullet a", "bullet b"])
    assert result["ok"] == "true", result
    edited = Path(result["edited_path"])
    assert edited != p
    assert edited.exists()
    # Original byte-for-byte unchanged.
    assert _sha256(p) == sha_orig
    # Edited copy has 2 slides.
    assert _pptx_slide_count(edited) == 2
    titles = _pptx_slide_titles(edited)
    assert titles[0] == "Cover"
    assert titles[1] == "Slide 2"


def test_add_slide_repeated_lands_in_same_edited(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p))
    r1 = office_add_slide.func(path=str(p), title="A")
    r2 = office_add_slide.func(path=str(p), title="B")
    assert r1["edited_path"] == r2["edited_path"]
    assert _pptx_slide_count(Path(r2["edited_path"])) == 3


def test_update_slide_changes_title_and_bullets(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="Initial")
    office_add_slide.func(path=str(p), title="Old title", bullets=["old"])
    edited = tmp_path / "deck.edited.pptx"

    result = office_update_slide.func(
        path=str(edited), index=1, title="New title", bullets=["one", "two", "three"]
    )
    assert result["ok"] == "true", result
    from pptx import Presentation

    prs = Presentation(edited)
    slide = prs.slides[1]
    assert slide.shapes.title.text_frame.text == "New title"
    # Bullets count check — find the body placeholder
    bodies = [s for s in slide.placeholders if s != slide.shapes.title]
    body_text = bodies[0].text_frame.text
    assert "one" in body_text
    assert "two" in body_text
    assert "three" in body_text


def test_delete_slide_removes_from_deck(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="First")
    office_add_slide.func(path=str(p), title="Second")
    office_add_slide.func(path=str(p), title="Third")
    edited = tmp_path / "deck.edited.pptx"
    assert _pptx_slide_count(edited) == 3

    result = office_delete_slide.func(path=str(edited), index=1)
    assert result["ok"] == "true", result
    titles = _pptx_slide_titles(edited)
    # "Second" deleted; First + Third remain in order
    assert titles == ["First", "Third"]


def test_reorder_slides_moves_position(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="A")
    office_add_slide.func(path=str(p), title="B")
    office_add_slide.func(path=str(p), title="C")
    edited = tmp_path / "deck.edited.pptx"
    # Move A (index 0) to the end (index 2).
    result = office_reorder_slides.func(path=str(edited), from_index=0, to_index=2)
    assert result["ok"] == "true"
    assert _pptx_slide_titles(edited) == ["B", "C", "A"]


def test_set_slide_title_convenience_tool(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="Old")
    edited_result = office_set_slide_title.func(path=str(p), index=0, title="Brand new")
    assert edited_result["ok"] == "true"
    edited = Path(edited_result["edited_path"])
    assert _pptx_slide_titles(edited)[0] == "Brand new"


def test_set_slide_bullets_convenience_tool(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p))
    office_add_slide.func(path=str(p), title="X", bullets=["old"])
    edited = tmp_path / "deck.edited.pptx"

    result = office_set_slide_bullets.func(path=str(edited), index=1, bullets=["new1", "new2"])
    assert result["ok"] == "true"
    from pptx import Presentation

    prs = Presentation(edited)
    bodies = [s for s in prs.slides[1].placeholders if s != prs.slides[1].shapes.title]
    body_text = bodies[0].text_frame.text
    assert "new1" in body_text and "new2" in body_text
    assert "old" not in body_text


def test_set_slide_notes_writes_speaker_notes(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p))
    result = office_set_slide_notes.func(path=str(p), index=0, notes="Remember to smile")
    assert result["ok"] == "true"
    edited = Path(result["edited_path"])
    from pptx import Presentation

    prs = Presentation(edited)
    assert "Remember to smile" in prs.slides[0].notes_slide.notes_text_frame.text


def test_add_image_to_slide_adds_picture_shape(tmp_path: Path) -> None:
    # Create a tiny PNG via stdlib so we don't add Pillow to test deps.
    # 1x1 transparent PNG bytes (well-known minimal).
    png_bytes = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
        "0000000a49444154789c6300010000000500010d0a2db40000000049454e44ae426082"
    )
    img = tmp_path / "tiny.png"
    img.write_bytes(png_bytes)
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p))

    result = office_add_image_to_slide.func(path=str(p), index=0, image_path=str(img))
    assert result["ok"] == "true", result
    edited = Path(result["edited_path"])
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(edited)
    # Slide should now contain at least one PICTURE shape (type 13).
    shapes = list(prs.slides[0].shapes)
    pictures = [sh for sh in shapes if sh.shape_type == 13]
    assert len(pictures) == 1
    # Default width should be 5 inches → 5 * 914400 EMU.
    assert abs(pictures[0].width - Emu(5 * 914400)) < 100


def test_read_slides_returns_structured_outline(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="Cover")
    office_add_slide.func(path=str(p), title="Findings", bullets=["A", "B"])
    office_set_slide_notes.func(path=str(p), index=1, notes="Talk slowly")
    edited = tmp_path / "deck.edited.pptx"

    result = office_read_slides.func(path=str(edited))
    assert result["ok"] == "true"
    assert result["slide_count"] == 2
    s0 = result["slides"][0]
    s1 = result["slides"][1]
    assert s0["title"] == "Cover"
    assert s1["title"] == "Findings"
    assert "A" in s1["bullets"]
    assert "B" in s1["bullets"]
    assert s1["notes"] == "Talk slowly"


def test_outline_dispatches_pptx_via_office_outline(tmp_path: Path) -> None:
    """office.outline (the read tool, not _outline_pptx helper) now
    also supports .pptx files."""
    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="Topic")
    result = office_outline.func(path=str(p))
    assert result["ok"] == "true"
    assert result["kind"] == "powerpoint"
    assert result["slide_count"] == 1
    assert result["slides"][0]["title"] == "Topic"


def test_pptx_atomic_rollback_on_verify_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Force-verify-fail on a pptx write → target keeps last-good state."""
    from shejane_runtime.tools import office as office_module

    p = tmp_path / "deck.pptx"
    office_create_pptx.func(path=str(p), title="A")
    # Do one real edit so .edited.pptx exists.
    r1 = office_add_slide.func(path=str(p), title="B")
    assert r1["ok"] == "true"
    target = Path(r1["edited_path"])
    sha_known_good = _sha256(target)

    monkeypatch.setattr(
        office_module,
        "_verify_file",
        lambda path, kind: (_ for _ in ()).throw(RuntimeError("boom")),
    )
    r2 = office_add_slide.func(path=str(p), title="C")
    assert r2["ok"] == "false"
    assert _sha256(target) == sha_known_good
