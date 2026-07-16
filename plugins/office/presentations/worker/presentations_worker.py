#!/usr/bin/env python3
"""One-shot Presentations Managed Worker for SheJane Plugin Action Protocol v1."""

from __future__ import annotations

import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from datetime import datetime
from pathlib import Path, PurePosixPath
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt

PLUGIN_ID = "org.shejane.presentations"
RUNTIME_ASSET_ID = "org.libreoffice.runtime"
ACTION_IDS = {
    "presentation.read",
    "presentation.create",
    "presentation.edit",
    "presentation.render",
}
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
IMAGE_MIMES = {"image/png", "image/jpeg", "image/gif"}
MAX_SLIDES = 200
MAX_TEXT = 60_000
MAX_TABLE_CELLS = 10_000


class PresentationActionError(RuntimeError):
    def __init__(self, code: str, message: str, *, retryable: bool = False) -> None:
        super().__init__(message)
        self.code = code
        self.retryable = retryable


def _reply(request_id: int, result: object) -> None:
    print(
        json.dumps(
            {"jsonrpc": "2.0", "id": request_id, "result": result},
            ensure_ascii=False,
            separators=(",", ":"),
        ),
        flush=True,
    )


def _request(expected_method: str) -> dict[str, Any]:
    try:
        payload = json.loads(sys.stdin.readline())
    except (EOFError, json.JSONDecodeError) as exc:
        raise PresentationActionError("protocol_violation", "invalid protocol frame") from exc
    if (
        not isinstance(payload, dict)
        or payload.get("jsonrpc") != "2.0"
        or payload.get("method") != expected_method
        or not isinstance(payload.get("params"), dict)
    ):
        raise PresentationActionError("protocol_violation", f"expected {expected_method}")
    return payload


def _output_root() -> Path:
    try:
        return Path(os.environ["SHEJANE_PLUGIN_OUTPUT_ROOT"]).resolve(strict=True)
    except (KeyError, OSError) as exc:
        raise PresentationActionError(
            "invalid_invocation", "output staging is unavailable"
        ) from exc


def _materialized_input(
    invocation: dict[str, Any], input_id: str, allowed_media_types: set[str]
) -> Path:
    references = invocation.get("inputs")
    if not isinstance(references, list):
        raise PresentationActionError("invalid_invocation", "authorized input is unavailable")
    reference = next(
        (item for item in references if isinstance(item, dict) and item.get("id") == input_id),
        None,
    )
    if reference is None or reference.get("media_type") not in allowed_media_types:
        raise PresentationActionError("invalid_invocation", "authorized input type is invalid")
    try:
        relative = PurePosixPath(str(reference["path"])).relative_to("/input")
    except (KeyError, ValueError) as exc:
        raise PresentationActionError("invalid_invocation", "input path is invalid") from exc
    if not relative.parts or any(part in {"", ".", ".."} for part in relative.parts):
        raise PresentationActionError("invalid_invocation", "input path is invalid")
    try:
        root = Path(os.environ["SHEJANE_PLUGIN_INPUT_ROOT"]).resolve(strict=True)
        source = root.joinpath(*relative.parts)
        if source.is_symlink() or not source.is_file():
            raise OSError
        source = source.resolve(strict=True)
        source.relative_to(root)
        data = source.read_bytes()
    except (KeyError, OSError, ValueError) as exc:
        raise PresentationActionError("invalid_invocation", "authorized input is unavailable") from exc
    if len(data) != reference.get("size_bytes") or hashlib.sha256(
        data
    ).hexdigest() != reference.get("sha256"):
        raise PresentationActionError("invalid_invocation", "input digest changed")
    return source


def _pptx_input(invocation: dict[str, Any], input_id: str) -> Path:
    source = _materialized_input(invocation, input_id, {PPTX_MIME})
    try:
        prefix = source.read_bytes()[:8]
    except OSError as exc:
        raise PresentationActionError("invalid_invocation", "deck input is unavailable") from exc
    if prefix.startswith(bytes.fromhex("d0cf11e0")):
        raise PresentationActionError(
            "presentation_encrypted_or_legacy",
            "encrypted or legacy binary PowerPoint decks are unsupported",
        )
    if not prefix.startswith(b"PK"):
        raise PresentationActionError("presentation_corrupt", "PPTX package is corrupt")
    try:
        Presentation(source)
    except Exception as exc:
        raise PresentationActionError("presentation_corrupt", "PPTX package is corrupt") from exc
    return source


def _safe_output_name(name: str) -> str:
    candidate = PurePosixPath(name)
    if (
        not name.lower().endswith(".pptx")
        or candidate.name != name
        or name in {".", ".."}
        or "\\" in name
    ):
        raise PresentationActionError("invalid_invocation", "output filename must be a basename.pptx")
    return name


def _bounded(text: Any, budget: list[int]) -> str:
    value = str(text or "")
    available = max(0, budget[0])
    result = value[:available]
    budget[0] -= len(result)
    return result


def _geometry(shape) -> dict[str, float]:
    return {
        "left": round(float(shape.left) / 914_400, 4),
        "top": round(float(shape.top) / 914_400, 4),
        "width": round(float(shape.width) / 914_400, 4),
        "height": round(float(shape.height) / 914_400, 4),
    }


def _title(slide) -> str:
    shape = slide.shapes.title
    return shape.text.strip() if shape is not None and shape.has_text_frame else ""


def _placeholder(slide, kinds: set[PP_PLACEHOLDER]):
    for shape in slide.placeholders:
        if shape.placeholder_format.type in kinds:
            return shape
    return None


def _bullets(slide, budget: list[int]) -> list[dict[str, Any]]:
    body = _placeholder(slide, {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT})
    if body is None or not body.has_text_frame:
        return []
    return [
        {"text": _bounded(paragraph.text, budget), "level": int(paragraph.level)}
        for paragraph in body.text_frame.paragraphs
        if paragraph.text.strip()
    ]


def _notes(slide, budget: list[int]) -> str:
    return (
        _bounded(slide.notes_slide.notes_text_frame.text.strip(), budget)
        if slide.has_notes_slide
        else ""
    )


def _unsupported_features(path: Path) -> list[str]:
    features: set[str] = set()
    try:
        with zipfile.ZipFile(path) as archive:
            for name in archive.namelist():
                lower = name.lower()
                if lower.startswith("ppt/slides/slide") and lower.endswith(".xml"):
                    data = archive.read(name)
                    if b"<p:timing" in data:
                        features.add("animations")
                    if b"<p:transition" in data:
                        features.add("transitions")
                if lower.startswith("ppt/media/") and lower.endswith(
                    (".mp3", ".wav", ".m4a", ".aac", ".mp4", ".mov", ".avi", ".wmv", ".webm")
                ):
                    features.add("audio_or_video")
    except (OSError, zipfile.BadZipFile) as exc:
        raise PresentationActionError("presentation_corrupt", "PPTX package is corrupt") from exc
    return sorted(features)


def _slide_details(slide, index: int, budget: list[int]) -> dict[str, Any]:
    title_shape = slide.shapes.title
    title_element = title_shape._element if title_shape is not None else None
    body = _placeholder(slide, {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT})
    body_element = body._element if body is not None else None
    subtitle = _placeholder(slide, {PP_PLACEHOLDER.SUBTITLE})
    subtitle_element = subtitle._element if subtitle is not None else None
    images: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    text_boxes: list[dict[str, Any]] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append(
                {
                    **_geometry(shape),
                    "media_type": shape.image.content_type,
                    "sha256": hashlib.sha256(shape.image.blob).hexdigest(),
                }
            )
        elif shape.has_table:
            rows = [
                [_bounded(cell.text, budget) for cell in row.cells]
                for row in shape.table.rows
            ]
            tables.append({**_geometry(shape), "rows": rows})
        elif (
            shape.has_text_frame
            and shape._element not in {title_element, body_element, subtitle_element}
            and shape.text.strip()
        ):
            text_boxes.append({**_geometry(shape), "text": _bounded(shape.text.strip(), budget)})
    return {
        "index": index,
        "layout": slide.slide_layout.name if slide.slide_layout is not None else "",
        "title": _bounded(_title(slide), budget),
        "subtitle": _bounded(subtitle.text.strip(), budget) if subtitle is not None else "",
        "bullets": _bullets(slide, budget),
        "notes": _notes(slide, budget),
        "image_count": len(images),
        "images": images,
        "tables": tables,
        "text_boxes": text_boxes,
        "shape_count": len(slide.shapes),
    }


def _presentation_read(
    invocation: dict[str, Any],
) -> tuple[dict[str, Any], list[dict[str, str]]]:
    arguments = invocation["arguments"]
    source = _pptx_input(invocation, str(arguments["input_id"]))
    deck = Presentation(source)
    start = max(0, int(arguments.get("start_slide", 0)))
    maximum = max(1, min(MAX_SLIDES, int(arguments.get("max_slides", 50))))
    end = min(len(deck.slides), start + maximum)
    budget = [MAX_TEXT]
    slides = [_slide_details(deck.slides[index], index, budget) for index in range(start, end)]
    truncated = end < len(deck.slides) or budget[0] == 0
    warnings = ["presentation output was truncated"] if truncated else []
    return (
        {
            "slide_count": len(deck.slides),
            "slide_width": round(float(deck.slide_width) / 914_400, 4),
            "slide_height": round(float(deck.slide_height) / 914_400, 4),
            "master_count": len(deck.slide_masters),
            "slides": slides,
            "unsupported_features": _unsupported_features(source),
            "truncated": truncated,
            "warnings": warnings,
        },
        [],
    )


def _layout(deck, name: str | None):
    requested = name or "Title and Content"
    for layout in deck.slide_layouts:
        if layout.name == requested:
            return layout
    raise PresentationActionError("presentation_layout_unsupported", "slide layout is unavailable")


def _inch(value: Any, *, positive: bool = False) -> int:
    try:
        numeric = float(value)
    except (TypeError, ValueError) as exc:
        raise PresentationActionError("invalid_invocation", "shape geometry is invalid") from exc
    if numeric < 0 or numeric > 100 or (positive and numeric <= 0):
        raise PresentationActionError("invalid_invocation", "shape geometry is invalid")
    return Inches(numeric)


def _set_text(shape, text: str, *, font_name: str | None = None, font_size: float | None = None) -> None:
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_size is not None:
        run.font.size = Pt(float(font_size))


def _set_bullets(
    shape, bullets: list[dict[str, Any]], *, default_font_name: str | None = None
) -> None:
    frame = shape.text_frame
    frame.clear()
    if not bullets:
        return
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(item["text"])
        paragraph.level = int(item.get("level", 0))
        for run in paragraph.runs:
            if item.get("font_name") or default_font_name:
                run.font.name = str(item.get("font_name") or default_font_name)
            if item.get("font_size") is not None:
                run.font.size = Pt(float(item["font_size"]))


def _add_text_box(slide, item: dict[str, Any]) -> None:
    shape = slide.shapes.add_textbox(
        _inch(item["left"]),
        _inch(item["top"]),
        _inch(item["width"], positive=True),
        _inch(item["height"], positive=True),
    )
    _set_text(
        shape,
        str(item["text"]),
        font_name=item.get("font_name"),
        font_size=item.get("font_size"),
    )


def _add_table(slide, item: dict[str, Any]) -> None:
    rows = item["rows"]
    width = len(rows[0])
    if not rows or not width or any(len(row) != width for row in rows):
        raise PresentationActionError("invalid_invocation", "table rows must have equal width")
    if len(rows) * width > MAX_TABLE_CELLS:
        raise PresentationActionError("invalid_invocation", "table exceeds the action cell limit")
    table = slide.shapes.add_table(
        len(rows),
        width,
        _inch(item["left"]),
        _inch(item["top"]),
        _inch(item["width"], positive=True),
        _inch(item["height"], positive=True),
    ).table
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = str(value)


def _add_image(slide, item: dict[str, Any], invocation: dict[str, Any]) -> None:
    source = _materialized_input(invocation, str(item["input_id"]), IMAGE_MIMES)
    kwargs: dict[str, Any] = {"width": _inch(item["width"], positive=True)}
    if item.get("height") is not None:
        kwargs["height"] = _inch(item["height"], positive=True)
    try:
        slide.shapes.add_picture(
            str(source),
            _inch(item["left"]),
            _inch(item["top"]),
            **kwargs,
        )
    except Exception as exc:
        raise PresentationActionError("invalid_invocation", "image input is invalid") from exc


def _populate_slide(
    deck,
    specification: dict[str, Any],
    invocation: dict[str, Any],
):
    slide = deck.slides.add_slide(_layout(deck, specification.get("layout")))
    if "title" in specification:
        if slide.shapes.title is None:
            raise PresentationActionError(
                "presentation_layout_unsupported", "slide layout has no title placeholder"
            )
        _set_text(
            slide.shapes.title,
            str(specification["title"]),
            font_name=specification.get("font_name"),
        )
    if "subtitle" in specification:
        subtitle = _placeholder(slide, {PP_PLACEHOLDER.SUBTITLE})
        if subtitle is None:
            raise PresentationActionError(
                "presentation_layout_unsupported", "slide layout has no subtitle placeholder"
            )
        _set_text(subtitle, str(specification["subtitle"]), font_name=specification.get("font_name"))
    if "bullets" in specification:
        body = _placeholder(slide, {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT})
        if body is None:
            raise PresentationActionError(
                "presentation_layout_unsupported", "slide layout has no content placeholder"
            )
        _set_bullets(
            body,
            specification["bullets"],
            default_font_name=specification.get("font_name"),
        )
    if "notes" in specification:
        slide.notes_slide.notes_text_frame.text = str(specification["notes"])
    for item in specification.get("text_boxes", []):
        _add_text_box(slide, item)
    for item in specification.get("tables", []):
        _add_table(slide, item)
    for item in specification.get("images", []):
        _add_image(slide, item, invocation)
    return slide


def _delete_slide(deck, index: int) -> None:
    if not 0 <= index < len(deck.slides):
        raise PresentationActionError("presentation_slide_not_found", "slide index is unavailable")
    slide_id = list(deck.slides._sldIdLst)[index]
    relationship = slide_id.attrib[
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    ]
    deck.part.drop_rel(relationship)
    deck.slides._sldIdLst.remove(slide_id)


def _move_slide(deck, from_index: int, to_index: int) -> None:
    slides = list(deck.slides._sldIdLst)
    if not 0 <= from_index < len(slides) or not 0 <= to_index < len(slides):
        raise PresentationActionError("presentation_slide_not_found", "slide index is unavailable")
    moving = slides[from_index]
    deck.slides._sldIdLst.remove(moving)
    deck.slides._sldIdLst.insert(to_index, moving)


def _normalize_archive(path: Path) -> None:
    normalized = path.with_suffix(".normalized.pptx")
    try:
        with (
            zipfile.ZipFile(path) as source,
            zipfile.ZipFile(normalized, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as target,
        ):
            for original in sorted(source.infolist(), key=lambda item: item.filename):
                if original.is_dir():
                    continue
                data = source.read(original)
                if original.filename == "docProps/core.xml":
                    data = re.sub(
                        rb"(<dcterms:(?:created|modified)\b[^>]*>)[^<]*(</dcterms:(?:created|modified)>)",
                        rb"\g<1>1980-01-01T00:00:00Z\g<2>",
                        data,
                    )
                info = zipfile.ZipInfo(original.filename, date_time=(1980, 1, 1, 0, 0, 0))
                info.compress_type = zipfile.ZIP_DEFLATED
                info.create_system = 3
                info.external_attr = 0o600 << 16
                target.writestr(info, data)
        os.replace(normalized, path)
    except (OSError, zipfile.BadZipFile) as exc:
        raise PresentationActionError(
            "presentation_write_failed", "PPTX output could not be normalized"
        ) from exc
    finally:
        normalized.unlink(missing_ok=True)


def _atomic_save(deck, filename: str) -> Path:
    output = _output_root()
    target = output / _safe_output_name(filename)
    if target.exists():
        raise PresentationActionError("presentation_write_failed", "PPTX output already exists")
    descriptor, temporary_name = tempfile.mkstemp(prefix=".pptx-", suffix=".pptx", dir=output)
    os.close(descriptor)
    temporary = Path(temporary_name)
    try:
        fixed = datetime(1980, 1, 1)
        deck.core_properties.created = fixed
        deck.core_properties.modified = fixed
        deck.core_properties.last_modified_by = "SheJane"
        deck.save(temporary)
        _normalize_archive(temporary)
        Presentation(temporary)
        os.replace(temporary, target)
        return target
    except PresentationActionError:
        raise
    except Exception as exc:
        raise PresentationActionError(
            "presentation_write_failed", "PPTX output could not be written"
        ) from exc
    finally:
        temporary.unlink(missing_ok=True)


def _presentation_create(
    invocation: dict[str, Any],
) -> tuple[dict[str, Any], list[dict[str, str]]]:
    arguments = invocation["arguments"]
    slides = arguments.get("slides")
    if not isinstance(slides, list) or not 1 <= len(slides) <= MAX_SLIDES:
        raise PresentationActionError("invalid_invocation", "one to 200 slides are required")
    if arguments.get("template_input_id"):
        template = _pptx_input(invocation, str(arguments["template_input_id"]))
        deck = Presentation(template)
        while deck.slides:
            _delete_slide(deck, 0)
    else:
        deck = Presentation()
        deck.slide_width = Inches(13.333333)
        deck.slide_height = Inches(7.5)
    for specification in slides:
        _populate_slide(deck, specification, invocation)
    target = _atomic_save(deck, str(arguments["filename"]))
    return (
        {
            "filename": target.name,
            "slide_count": len(deck.slides),
            "master_count": len(deck.slide_masters),
            "warnings": [],
        },
        [{"path": f"/output/{target.name}", "media_type": PPTX_MIME, "name": target.name}],
    )


def _update_slide(slide, operation: dict[str, Any]) -> None:
    if "title" in operation:
        if slide.shapes.title is None:
            raise PresentationActionError(
                "presentation_layout_unsupported", "slide layout has no title placeholder"
            )
        _set_text(slide.shapes.title, str(operation["title"]))
    if "subtitle" in operation:
        subtitle = _placeholder(slide, {PP_PLACEHOLDER.SUBTITLE})
        if subtitle is None:
            raise PresentationActionError(
                "presentation_layout_unsupported", "slide layout has no subtitle placeholder"
            )
        _set_text(subtitle, str(operation["subtitle"]))
    if "bullets" in operation:
        body = _placeholder(slide, {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT})
        if body is None:
            raise PresentationActionError(
                "presentation_layout_unsupported", "slide layout has no content placeholder"
            )
        _set_bullets(body, operation["bullets"])
    if "notes" in operation:
        slide.notes_slide.notes_text_frame.text = str(operation["notes"])


def _presentation_edit(
    invocation: dict[str, Any],
) -> tuple[dict[str, Any], list[dict[str, str]]]:
    arguments = invocation["arguments"]
    source = _pptx_input(invocation, str(arguments["input_id"]))
    operations = arguments.get("operations")
    if not isinstance(operations, list) or not 1 <= len(operations) <= 128:
        raise PresentationActionError("invalid_invocation", "one to 128 operations are required")
    deck = Presentation(source)
    changes: list[str] = []
    for operation in operations:
        operation_type = operation["type"]
        if operation_type == "add_slide":
            if len(deck.slides) >= MAX_SLIDES:
                raise PresentationActionError("invalid_invocation", "deck exceeds the slide limit")
            _populate_slide(deck, operation["slide"], invocation)
            changes.append("add_slide:1")
        elif operation_type == "update_slide":
            index = int(operation["index"])
            if not 0 <= index < len(deck.slides):
                raise PresentationActionError(
                    "presentation_slide_not_found", "slide index is unavailable"
                )
            _update_slide(deck.slides[index], operation["content"])
            changes.append("update_slide:1")
        elif operation_type == "delete_slide":
            if len(deck.slides) == 1:
                raise PresentationActionError("invalid_invocation", "the last slide cannot be deleted")
            _delete_slide(deck, int(operation["index"]))
            changes.append("delete_slide:1")
        elif operation_type == "move_slide":
            _move_slide(deck, int(operation["from_index"]), int(operation["to_index"]))
            changes.append("move_slide:1")
        elif operation_type in {"add_image", "add_text_box", "add_table"}:
            index = int(operation["index"])
            if not 0 <= index < len(deck.slides):
                raise PresentationActionError(
                    "presentation_slide_not_found", "slide index is unavailable"
                )
            item = operation[operation_type.removeprefix("add_")]
            if operation_type == "add_image":
                _add_image(deck.slides[index], item, invocation)
            elif operation_type == "add_text_box":
                _add_text_box(deck.slides[index], item)
            else:
                _add_table(deck.slides[index], item)
            changes.append(f"{operation_type}:1")
        else:
            raise PresentationActionError(
                "invalid_invocation", "presentation edit operation is unsupported"
            )
    target = _atomic_save(
        deck,
        str(arguments.get("output_filename") or f"{source.stem}.edited.pptx"),
    )
    return (
        {
            "filename": target.name,
            "slide_count": len(deck.slides),
            "operation_count": len(operations),
            "changes": changes,
            "warnings": [],
        },
        [{"path": f"/output/{target.name}", "media_type": PPTX_MIME, "name": target.name}],
    )


def _runtime_tool_paths() -> tuple[Path, Path]:
    try:
        mapping = json.loads(os.environ["SHEJANE_PLUGIN_RUNTIME_ASSETS"])
        root = Path(mapping[RUNTIME_ASSET_ID]).resolve(strict=True)
        config = json.loads((root / "office-runtime.json").read_text(encoding="utf-8"))
    except (KeyError, OSError, ValueError, json.JSONDecodeError) as exc:
        raise PresentationActionError(
            "engine_unavailable", "Office runtime asset is unavailable"
        ) from exc
    tools = []
    for key in ("soffice", "mutool"):
        try:
            relative = PurePosixPath(str(config[key]))
            if relative.is_absolute() or any(part in {"", ".", ".."} for part in relative.parts):
                raise ValueError
            candidate = root.joinpath(*relative.parts).resolve(strict=True)
            candidate.relative_to(root)
        except (KeyError, OSError, ValueError) as exc:
            raise PresentationActionError(
                "engine_unavailable", "Office runtime asset is invalid"
            ) from exc
        if not candidate.is_file() or not os.access(candidate, os.X_OK):
            raise PresentationActionError(
                "engine_unavailable", "Office runtime executable is unavailable"
            )
        tools.append(candidate)
    return tools[0], tools[1]


def _run_tool(command: list[str], *, cwd: Path, deadline: float, env: dict[str, str]) -> str:
    remaining = deadline - time.monotonic()
    if remaining <= 0:
        raise PresentationActionError("render_failed", "presentation renderer timed out")
    try:
        completed = subprocess.run(
            command,
            cwd=cwd,
            env=env,
            stdin=subprocess.DEVNULL,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            timeout=remaining,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        raise PresentationActionError(
            "render_failed", "presentation renderer did not complete"
        ) from exc
    if completed.returncode != 0:
        raise PresentationActionError("render_failed", "presentation renderer rejected the input")
    return completed.stdout[-16_384:].decode("utf-8", errors="replace")


def _presentation_render(
    invocation: dict[str, Any],
) -> tuple[dict[str, Any], list[dict[str, str]]]:
    arguments = invocation["arguments"]
    source = _pptx_input(invocation, str(arguments["input_id"]))
    deck = Presentation(source)
    unsupported = _unsupported_features(source)
    soffice, mutool = _runtime_tool_paths()
    output = _output_root()
    profile = output / ".libreoffice-profile"
    temporary = output / ".runtime-tmp"
    profile.mkdir(mode=0o700, exist_ok=False)
    temporary.mkdir(mode=0o700, exist_ok=False)
    deadline = time.monotonic() + max(
        1.0, min(300.0, float(invocation["limits"]["timeout_ms"]) / 1000 * 0.8)
    )
    env = {
        "HOME": str(profile),
        "TMPDIR": str(temporary),
        "TEMP": str(temporary),
        "TMP": str(temporary),
        "LANG": "C.UTF-8",
        "LC_ALL": "C.UTF-8",
        "TZ": "UTC",
        "PATH": os.pathsep.join(sorted({str(soffice.parent), str(mutool.parent)})),
    }
    try:
        _run_tool(
            [
                str(soffice),
                "--headless",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--norestore",
                f"-env:UserInstallation={profile.as_uri()}",
                "--convert-to",
                "pdf:impress_pdf_Export",
                "--outdir",
                str(output),
                str(source),
            ],
            cwd=output,
            deadline=deadline,
            env=env,
        )
        pdf = output / f"{source.stem}.pdf"
        if not pdf.is_file() or not pdf.read_bytes().startswith(b"%PDF"):
            raise PresentationActionError(
                "render_failed", "presentation renderer produced no valid PDF"
            )
        info = _run_tool(
            [str(mutool), "info", str(pdf)], cwd=output, deadline=deadline, env=env
        )
        match = re.search(r"(?:Pages|pages)\s*:\s*(\d+)", info)
        if match is None or int(match.group(1)) < 1:
            raise PresentationActionError("render_failed", "PDF page count is unavailable")
        page_count = int(match.group(1))
        include_png = bool(arguments.get("include_png", True))
        rendered = min(page_count, int(arguments.get("max_slides", 20))) if include_png else 0
        png_names: list[str] = []
        if rendered:
            pattern = output / f"{source.stem}.slide-%04d.png"
            _run_tool(
                [
                    str(mutool),
                    "draw",
                    "-q",
                    "-r",
                    str(int(arguments.get("dpi", 144))),
                    "-o",
                    str(pattern),
                    str(pdf),
                    f"1-{rendered}",
                ],
                cwd=output,
                deadline=deadline,
                env=env,
            )
            for page in range(1, rendered + 1):
                name = f"{source.stem}.slide-{page:04d}.png"
                if not (output / name).is_file():
                    raise PresentationActionError(
                        "render_failed", "PNG slide rendering was incomplete"
                    )
                png_names.append(name)
        warnings = [f"{feature} are not represented in PDF preview" for feature in unsupported]
        if include_png and rendered < page_count:
            warnings.append("PNG preview limited by max_slides")
        if page_count != len(deck.slides):
            warnings.append("PDF page count differs from deck slide count")
        return (
            {
                "slide_count": len(deck.slides),
                "rendered_slides": rendered,
                "pdf_name": pdf.name,
                "png_names": png_names,
                "unsupported_features": unsupported,
                "warnings": warnings,
            },
            [
                {"path": f"/output/{pdf.name}", "media_type": "application/pdf", "name": pdf.name},
                *(
                    {"path": f"/output/{name}", "media_type": "image/png", "name": name}
                    for name in png_names
                ),
            ],
        )
    finally:
        shutil.rmtree(profile, ignore_errors=True)
        shutil.rmtree(temporary, ignore_errors=True)


def _success(
    invocation: dict[str, Any], output: dict[str, Any], artifacts: list[dict[str, str]]
) -> dict[str, Any]:
    return {
        "schema_version": 1,
        "invocation_id": invocation["invocation_id"],
        "operation_id": invocation["operation_id"],
        "status": "succeeded",
        "output": output,
        "artifacts": artifacts,
    }


def _failure(invocation: dict[str, Any], error: PresentationActionError) -> dict[str, Any]:
    return {
        "schema_version": 1,
        "invocation_id": invocation["invocation_id"],
        "operation_id": invocation["operation_id"],
        "status": "failed",
        "error": {"code": error.code, "message": str(error), "retryable": error.retryable},
        "artifacts": [],
    }


def _invoke(invocation: dict[str, Any]) -> dict[str, Any]:
    action = invocation.get("action")
    if not isinstance(action, dict) or action.get("plugin_id") != PLUGIN_ID:
        raise PresentationActionError("invalid_invocation", "plugin identity does not match")
    handlers = {
        "presentation.read": _presentation_read,
        "presentation.create": _presentation_create,
        "presentation.edit": _presentation_edit,
        "presentation.render": _presentation_render,
    }
    handler = handlers.get(str(action.get("action_id") or ""))
    if handler is None:
        raise PresentationActionError("invalid_invocation", "presentation Action is unsupported")
    output, artifacts = handler(invocation)
    return _success(invocation, output, artifacts)


def main() -> None:
    initialize = _request("initialize")
    params = initialize["params"]
    if (
        params.get("protocol_version") != 1
        or params.get("plugin_id") != PLUGIN_ID
        or not set(params.get("actions", ())).issubset(ACTION_IDS)
    ):
        raise PresentationActionError("incompatible_runtime", "worker initialization is incompatible")
    _reply(
        int(initialize["id"]),
        {
            "protocol_version": 1,
            "process_isolated": True,
            "access_isolated": os.environ.get("SHEJANE_PLUGIN_ACCESS_ISOLATED") == "1",
            "resource_isolated": os.environ.get("SHEJANE_PLUGIN_RESOURCE_ISOLATED") == "1",
            "sandboxed": os.environ.get("SHEJANE_PLUGIN_SANDBOXED") == "1",
        },
    )
    invoke = _request("invoke")
    invocation = invoke["params"]
    try:
        result = _invoke(invocation)
    except PresentationActionError as exc:
        result = _failure(invocation, exc)
    except Exception as exc:
        if os.environ.get("SHEJANE_PLUGIN_DEBUG") == "1":
            print(f"{type(exc).__name__}: {exc}", file=sys.stderr, flush=True)
        result = _failure(
            invocation,
            PresentationActionError("plugin_failed", "presentation worker failed"),
        )
    _reply(int(invoke["id"]), result)
    shutdown = _request("shutdown")
    _reply(int(shutdown["id"]), {})


if __name__ == "__main__":
    main()
